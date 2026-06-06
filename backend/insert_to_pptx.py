from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path
from typing import Any

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Emu, Inches, Pt
except ModuleNotFoundError:  # pragma: no cover
    Presentation = None
    RGBColor = None
    PP_ALIGN = None
    Emu = None
    Inches = None
    Pt = None


@dataclass
class InsertResult:
    output_path: str
    slide_number: int
    chart_left: int
    chart_top: int
    chart_width: int
    chart_height: int
    replaced_table: bool

    def to_dict(self) -> dict[str, Any]:
        return {
            "output_path": self.output_path,
            "slide_number": self.slide_number,
            "chart_left": self.chart_left,
            "chart_top": self.chart_top,
            "chart_width": self.chart_width,
            "chart_height": self.chart_height,
            "replaced_table": self.replaced_table,
        }


def _ensure_dependencies() -> None:
    if Presentation is None:
        raise ModuleNotFoundError("python-pptx is required for PPT insertion.")


def _derive_output_path(ppt_path: str | Path, output_path: str | Path | None) -> Path:
    source = Path(ppt_path)
    return Path(output_path) if output_path is not None else source.with_name(f"{source.stem}_enhanced{source.suffix}")


def _is_raster_image(path: str | Path) -> bool:
    return Path(path).suffix.lower() in {".png", ".jpg", ".jpeg"}


def _shape_bounds(shape: dict[str, Any]) -> tuple[int, int, int, int]:
    left = int(shape.get("left", 0))
    top = int(shape.get("top", 0))
    width = int(shape.get("width", 0))
    height = int(shape.get("height", 0))
    return left, top, width, height


def _best_anchor(shapes: list[dict[str, Any]]) -> dict[str, int]:
    table_shapes = [shape for shape in shapes if shape.get("has_table")]
    if table_shapes:
        largest = max(table_shapes, key=lambda item: item.get("width", 0) * item.get("height", 0))
        return {
            "index": int(largest.get("index", 0)),
            "left": int(largest.get("left", 0)),
            "top": int(largest.get("top", 0)),
            "width": int(largest.get("width", 0)),
            "height": int(largest.get("height", 0)),
        }
    return {"index": 0, "left": int(Inches(0.8)), "top": int(Inches(1.5)), "width": int(Inches(8.0)), "height": int(Inches(4.5))}


def _horizontal_overlap(first: dict[str, Any], second: dict[str, Any]) -> int:
    first_left, _, first_width, _ = _shape_bounds(first)
    second_left, _, second_width, _ = _shape_bounds(second)
    first_right = first_left + first_width
    second_right = second_left + second_width
    return max(0, min(first_right, second_right) - max(first_left, second_left))


def _find_adjacent_text_shape(shapes: list[dict[str, Any]], anchor: dict[str, int], direction: str) -> dict[str, Any] | None:
    anchor_box = {
        "left": anchor.get("left", 0),
        "top": anchor.get("top", 0),
        "width": anchor.get("width", 0),
        "height": anchor.get("height", 0),
    }
    _, anchor_top, anchor_width, anchor_height = _shape_bounds(anchor_box)
    anchor_bottom = anchor_top + anchor_height
    min_overlap = max(int(anchor_width * 0.35), int(Inches(1.6)))
    max_gap = int(Inches(1.0))
    candidates: list[tuple[int, int, dict[str, Any]]] = []

    for shape in shapes:
        if not shape.get("has_text") or shape.get("has_table"):
            continue
        overlap = _horizontal_overlap(shape, anchor_box)
        if overlap < min_overlap:
            continue
        _, top, _, height = _shape_bounds(shape)
        bottom = top + height
        gap = anchor_top - bottom if direction == "above" else top - anchor_bottom
        if gap < 0 or gap > max_gap:
            continue
        candidates.append((gap, -overlap, shape))

    if not candidates:
        return None
    return min(candidates, key=lambda item: (item[0], item[1]))[2]


def _remove_shape(shape: Any) -> None:
    element = shape._element
    element.getparent().remove(element)


def _add_blank_slide(presentation: Any) -> Any:
    layout_count = len(presentation.slide_layouts)
    if layout_count <= 0:
        raise ValueError("Presentation does not contain any slide layouts.")
    layout_index = 6 if layout_count > 6 else layout_count - 1
    slide = presentation.slides.add_slide(presentation.slide_layouts[layout_index])
    for shape in list(slide.shapes):
        _remove_shape(shape)
    return slide


def _replace_table_shape(slide: Any, anchor_index: int) -> bool:
    if anchor_index <= 0:
        return False
    shapes = list(slide.shapes)
    if anchor_index > len(shapes):
        return False
    candidate = shapes[anchor_index - 1]
    if not getattr(candidate, "has_table", False):
        return False
    _remove_shape(candidate)
    return True


def _remove_shape_indexes(slide: Any, shape_indexes: list[int]) -> None:
    for shape_index in sorted({index for index in shape_indexes if index > 0}, reverse=True):
        shapes = list(slide.shapes)
        if shape_index > len(shapes):
            continue
        _remove_shape(shapes[shape_index - 1])


def _replace_table_region(slide: Any) -> tuple[int, int, int, int] | None:
    for shape in slide.shapes:
        if getattr(shape, "has_table", False):
            return int(shape.left), int(shape.top), int(shape.width), int(shape.height)
    return None


def _default_chart_region(slide_width: int, slide_height: int) -> tuple[int, int, int, int]:
    return (int(slide_width * 0.08), int(slide_height * 0.20), int(slide_width * 0.50), int(slide_height * 0.58))


def _default_illustration_region(slide_width: int, slide_height: int) -> tuple[int, int, int, int]:
    return (int(slide_width * 0.62), int(slide_height * 0.23), int(slide_width * 0.27), int(slide_height * 0.36))


def _result_slide_chart_region(slide_width: int, slide_height: int) -> tuple[int, int, int, int]:
    return (int(slide_width * 0.08), int(slide_height * 0.24), int(slide_width * 0.54), int(slide_height * 0.52))


def _result_slide_illustration_region(slide_width: int, slide_height: int) -> tuple[int, int, int, int]:
    return (int(slide_width * 0.68), int(slide_height * 0.28), int(slide_width * 0.24), int(slide_height * 0.34))


def _rect_area(region: tuple[int, int, int, int]) -> int:
    return max(0, region[2]) * max(0, region[3])


def _rect_to_edges(region: tuple[int, int, int, int]) -> tuple[int, int, int, int]:
    left, top, width, height = region
    return left, top, left + width, top + height


def _overlap_area(first: tuple[int, int, int, int], second: tuple[int, int, int, int]) -> int:
    first_edges = _rect_to_edges(first)
    second_edges = _rect_to_edges(second)
    left = max(first_edges[0], second_edges[0])
    top = max(first_edges[1], second_edges[1])
    right = min(first_edges[2], second_edges[2])
    bottom = min(first_edges[3], second_edges[3])
    return max(0, right - left) * max(0, bottom - top)


def _shape_region(shape: dict[str, Any]) -> tuple[int, int, int, int]:
    return (
        int(shape.get("left", 0)),
        int(shape.get("top", 0)),
        int(shape.get("width", 0)),
        int(shape.get("height", 0)),
    )


def _shape_overlap_weight(shape: dict[str, Any]) -> float:
    if shape.get("has_table"):
        return 4.0
    if shape.get("has_picture") or str(shape.get("shape_type", "")).upper() == "PICTURE":
        return 6.0
    if shape.get("has_text"):
        return 4.5
    return 0.75


def _overlap_score(region: tuple[int, int, int, int], shapes: list[dict[str, Any]], slide_area: int, ignored_shape_index: int | None = None) -> float:
    region_area = max(1, _rect_area(region))
    score = 0.0
    for shape in shapes:
        if ignored_shape_index is not None and int(shape.get("index", 0)) == ignored_shape_index:
            continue
        shape_region = _shape_region(shape)
        shape_area = _rect_area(shape_region)
        if shape_area <= 0:
            continue
        if shape_area >= slide_area * 0.92:
            continue
        overlap = _overlap_area(region, shape_region)
        if overlap <= 0:
            continue
        score += _shape_overlap_weight(shape) * (overlap / max(1, min(region_area, shape_area)))
    return score


def _bounded_region(slide_width: int, slide_height: int, left_ratio: float, top_ratio: float, width_ratio: float, height_ratio: float) -> tuple[int, int, int, int]:
    width = int(slide_width * width_ratio)
    height = int(slide_height * height_ratio)
    margin_x = int(slide_width * 0.04)
    margin_y = int(slide_height * 0.06)
    left = min(max(int(slide_width * left_ratio), margin_x), max(margin_x, slide_width - width - margin_x))
    top = min(max(int(slide_height * top_ratio), margin_y), max(margin_y, slide_height - height - margin_y))
    return left, top, width, height


def _dedupe_regions(regions: list[tuple[int, int, int, int]]) -> list[tuple[int, int, int, int]]:
    seen: set[tuple[int, int, int, int]] = set()
    deduped: list[tuple[int, int, int, int]] = []
    for region in regions:
        if region in seen:
            continue
        seen.add(region)
        deduped.append(region)
    return deduped


def _chart_region_candidates(slide_width: int, slide_height: int) -> list[tuple[int, int, int, int]]:
    regions = [_default_chart_region(slide_width, slide_height)]
    sizes = [(0.50, 0.42), (0.42, 0.34), (0.36, 0.28), (0.30, 0.24)]
    positions = [(0.08, 0.20), (0.38, 0.20), (0.56, 0.20), (0.08, 0.36), (0.38, 0.36), (0.56, 0.36), (0.08, 0.54), (0.38, 0.54), (0.56, 0.54), (0.08, 0.66), (0.38, 0.66), (0.56, 0.66)]
    for width_ratio, height_ratio in sizes:
        for left_ratio, top_ratio in positions:
            regions.append(_bounded_region(slide_width, slide_height, left_ratio, top_ratio, width_ratio, height_ratio))
    return _dedupe_regions(regions)


def _illustration_region_candidates(slide_width: int, slide_height: int) -> list[tuple[int, int, int, int]]:
    regions = [_default_illustration_region(slide_width, slide_height)]
    sizes = [(0.27, 0.30), (0.22, 0.24), (0.18, 0.20)]
    positions = [(0.08, 0.20), (0.34, 0.20), (0.62, 0.20), (0.08, 0.40), (0.34, 0.40), (0.62, 0.40), (0.08, 0.58), (0.34, 0.58), (0.62, 0.58), (0.08, 0.70), (0.34, 0.70), (0.62, 0.70)]
    for width_ratio, height_ratio in sizes:
        for left_ratio, top_ratio in positions:
            regions.append(_bounded_region(slide_width, slide_height, left_ratio, top_ratio, width_ratio, height_ratio))
    return _dedupe_regions(regions)


def _candidate_layouts(slide_width: int, slide_height: int) -> list[tuple[tuple[int, int, int, int], tuple[int, int, int, int]]]:
    fixed_layouts = [
        (
            _default_chart_region(slide_width, slide_height),
            _default_illustration_region(slide_width, slide_height),
        ),
        (
            (int(slide_width * 0.38), int(slide_height * 0.20), int(slide_width * 0.50), int(slide_height * 0.54)),
            (int(slide_width * 0.08), int(slide_height * 0.24), int(slide_width * 0.27), int(slide_height * 0.34)),
        ),
        (
            (int(slide_width * 0.08), int(slide_height * 0.50), int(slide_width * 0.50), int(slide_height * 0.38)),
            (int(slide_width * 0.62), int(slide_height * 0.53), int(slide_width * 0.27), int(slide_height * 0.30)),
        ),
        (
            (int(slide_width * 0.08), int(slide_height * 0.28), int(slide_width * 0.46), int(slide_height * 0.36)),
            (int(slide_width * 0.62), int(slide_height * 0.28), int(slide_width * 0.27), int(slide_height * 0.30)),
        ),
    ]
    layouts = list(fixed_layouts)
    for chart_region in _chart_region_candidates(slide_width, slide_height):
        for illustration_region in _illustration_region_candidates(slide_width, slide_height):
            if _overlap_area(chart_region, illustration_region) > 0:
                continue
            layouts.append((chart_region, illustration_region))
    return _dedupe_layouts(layouts)


def _dedupe_layouts(layouts: list[tuple[tuple[int, int, int, int], tuple[int, int, int, int]]]) -> list[tuple[tuple[int, int, int, int], tuple[int, int, int, int]]]:
    seen: set[tuple[tuple[int, int, int, int], tuple[int, int, int, int]]] = set()
    deduped: list[tuple[tuple[int, int, int, int], tuple[int, int, int, int]]] = []
    for layout in layouts:
        if layout in seen:
            continue
        seen.add(layout)
        deduped.append(layout)
    return deduped


def _choose_asset_regions(
    slide_width: int,
    slide_height: int,
    shapes: list[dict[str, Any]],
    chart_region: tuple[int, int, int, int] | None = None,
    ignored_shape_index: int | None = None,
) -> tuple[tuple[int, int, int, int], tuple[int, int, int, int]]:
    slide_area = max(1, slide_width * slide_height)
    if chart_region is not None:
        chart_candidates = [chart_region]
        illustration_candidates = _illustration_region_candidates(slide_width, slide_height)
    else:
        layouts = _candidate_layouts(slide_width, slide_height)
        best_layout = min(
            layouts,
            key=lambda item: (
                _overlap_score(item[0], shapes, slide_area, ignored_shape_index)
                + _overlap_score(item[1], shapes, slide_area, ignored_shape_index)
                + (_overlap_area(item[0], item[1]) / slide_area) * 100,
                _overlap_score(item[0], shapes, slide_area, ignored_shape_index),
            ),
        )
        return best_layout

    best_chart = chart_candidates[0]
    best_illustration = min(
        illustration_candidates,
        key=lambda item: (
            _overlap_score(item, shapes, slide_area, ignored_shape_index)
            + (_overlap_area(best_chart, item) / slide_area) * 100,
            _overlap_area(best_chart, item),
        ),
    )
    return best_chart, best_illustration


def _layout_metadata(
    chart_region: tuple[int, int, int, int],
    illustration_region: tuple[int, int, int, int],
    shapes: list[dict[str, Any]],
    slide_width: int,
    slide_height: int,
    ignored_shape_index: int | None = None,
) -> dict[str, Any]:
    slide_area = max(1, slide_width * slide_height)
    chart_score = _overlap_score(chart_region, shapes, slide_area, ignored_shape_index)
    illustration_score = _overlap_score(illustration_region, shapes, slide_area, ignored_shape_index)
    pair_overlap = _overlap_area(chart_region, illustration_region)
    overlap_score = round(chart_score + illustration_score + (pair_overlap / slide_area) * 100, 4)
    return {
        "chart_region": {
            "left": chart_region[0],
            "top": chart_region[1],
            "width": chart_region[2],
            "height": chart_region[3],
        },
        "illustration_region": {
            "left": illustration_region[0],
            "top": illustration_region[1],
            "width": illustration_region[2],
            "height": illustration_region[3],
        },
        "overlap_score": overlap_score,
        "layout_warning": overlap_score > 0.2,
        "pair_overlap": pair_overlap,
        "insertion_mode": "inline",
        "original_slide_preserved": False,
    }


def _compute_chart_region(
    base_region: tuple[int, int, int, int],
    slide_height: int,
    title_region: tuple[int, int, int, int] | None = None,
    legend_region: tuple[int, int, int, int] | None = None,
) -> tuple[int, int, int, int]:
    left, top, width, height = base_region
    bottom = top + height
    padding = int(Inches(0.12))
    min_height = int(Inches(1.8))

    if title_region is not None:
        title_bottom = title_region[1] + title_region[3]
        top = max(top, title_bottom + padding)
    if legend_region is not None:
        bottom = min(bottom, legend_region[1] - padding)

    if bottom - top < min_height:
        bottom = min(slide_height - padding, top + max(min_height, height))

    return left, top, width, max(min_height, bottom - top)


def _region_from_layout_override(
    slide_width: int,
    slide_height: int,
    override: dict[str, Any],
    prefix: str,
) -> tuple[int, int, int, int] | None:
    x_key = f"{prefix}X"
    y_key = f"{prefix}Y"
    scale_key = f"{prefix}Scale"
    if x_key not in override or y_key not in override or scale_key not in override:
        return None
    try:
        left_ratio = float(override[x_key]) / 100
        top_ratio = float(override[y_key]) / 100
        width_ratio = float(override[scale_key]) / 100
    except (TypeError, ValueError):
        return None
    width_ratio = min(max(width_ratio, 0.12), 0.86)
    width = int(slide_width * width_ratio)
    height = int(width * 9 / 16)
    margin_x = int(slide_width * 0.02)
    margin_y = int(slide_height * 0.03)
    left = min(max(int(slide_width * left_ratio), margin_x), max(margin_x, slide_width - width - margin_x))
    top = min(max(int(slide_height * top_ratio), margin_y), max(margin_y, slide_height - height - margin_y))
    return left, top, width, height


def _manual_override_regions(
    slide_width: int,
    slide_height: int,
    layout_override: dict[str, Any] | None,
) -> tuple[tuple[int, int, int, int], tuple[int, int, int, int]] | None:
    if not layout_override:
        return None
    chart_region = _region_from_layout_override(slide_width, slide_height, layout_override, "chart")
    illustration_region = _region_from_layout_override(slide_width, slide_height, layout_override, "illustration")
    if chart_region is None or illustration_region is None:
        return None
    return chart_region, illustration_region


def _add_asset_result_content(
    slide: Any,
    slide_width: int,
    slide_height: int,
    chart_region: tuple[int, int, int, int],
    illustration_region: tuple[int, int, int, int],
    chart_path: str | Path | None,
    illustration_path: str | Path | None,
    title: str,
    summary: str,
    intent: dict[str, Any],
) -> None:
    chart_type = intent.get("chart_type", "bar")
    model = intent.get("image_model", "local")
    style = intent.get("illustration_style", "auto")
    score = intent.get("clip_score")

    _add_textbox(slide, int(slide_width * 0.08), int(slide_height * 0.08), int(slide_width * 0.76), int(slide_height * 0.07), title, 24, bold=True)
    _add_caption_box(slide, int(slide_width * 0.08), int(slide_height * 0.15), int(slide_width * 0.76), int(slide_height * 0.05), summary)

    if chart_path and _is_raster_image(chart_path) and Path(chart_path).exists():
        slide.shapes.add_picture(str(chart_path), Emu(chart_region[0]), Emu(chart_region[1]), width=Emu(chart_region[2]), height=Emu(chart_region[3]))
    else:
        _add_caption_box(slide, chart_region[0], chart_region[1], chart_region[2], int(slide_height * 0.12), "Chart preview is not available as a raster image yet.")

    legend_text = f"Chart type: {chart_type} | Source: {intent.get('source', 'local')} | Mode: {intent.get('semantic_mode', 'local')}"
    _add_caption_box(slide, chart_region[0], chart_region[1] + chart_region[3] + int(slide_height * 0.02), chart_region[2], int(slide_height * 0.08), legend_text)

    if illustration_path and _is_raster_image(illustration_path) and Path(illustration_path).exists():
        slide.shapes.add_picture(str(illustration_path), Emu(illustration_region[0]), Emu(illustration_region[1]), width=Emu(illustration_region[2]), height=Emu(illustration_region[3]))
    else:
        _add_textbox(slide, illustration_region[0], illustration_region[1], illustration_region[2], int(slide_height * 0.08), "Illustration Area", 20, bold=True, color=(52, 82, 126))
        _add_caption_box(slide, illustration_region[0], illustration_region[1] + int(slide_height * 0.08), illustration_region[2], int(slide_height * 0.18), f"Style: {style} | Model: {model}\nMatch score: {score if score is not None else 'N/A'}\nTheme: {intent.get('visual_theme', 'Generated illustration preview')}")


def _add_textbox(slide: Any, left: int, top: int, width: int, height: int, text: str, font_size: int, bold: bool = False, color: tuple[int, int, int] = (32, 52, 82)) -> None:
    box = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    paragraph = frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.LEFT
    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor(*color)


def _add_title(slide: Any, title: str, region: tuple[int, int, int, int] | None = None) -> None:
    if region is None:
        region = (int(Inches(0.6)), int(Inches(0.25)), int(Inches(8.0)), int(Inches(0.55)))
    _add_textbox(slide, region[0], region[1], region[2], region[3], title, 22, bold=True, color=(35, 35, 35))


def _add_caption_box(slide: Any, left: int, top: int, width: int, height: int, text: str) -> None:
    box = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    paragraph = frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(14)
    run.font.color.rgb = RGBColor(88, 104, 127)


def _add_title_block(slide: Any, title: str, subtitle: str, region: tuple[int, int, int, int] | None, slide_width: int, slide_height: int) -> None:
    if region is None:
        title_height = int(slide_height * 0.07)
        subtitle_height = int(slide_height * 0.05)
        title_region = (int(slide_width * 0.08), int(slide_height * 0.08), int(slide_width * 0.76), title_height)
        _add_title(slide, title, title_region)
        if subtitle:
            subtitle_region = (int(slide_width * 0.08), int(slide_height * 0.15), int(slide_width * 0.76), subtitle_height)
            _add_caption_box(slide, *subtitle_region, subtitle)
        return

    left, top, width, height = region
    box = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    paragraph = frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.LEFT
    run = paragraph.add_run()
    run.text = title
    run.font.size = Pt(20)
    run.font.bold = True
    run.font.color.rgb = RGBColor(35, 35, 35)
    if subtitle:
        subtitle_paragraph = frame.add_paragraph()
        subtitle_paragraph.alignment = PP_ALIGN.LEFT
        subtitle_run = subtitle_paragraph.add_run()
        subtitle_run.text = subtitle
        subtitle_run.font.size = Pt(12)
        subtitle_run.font.color.rgb = RGBColor(88, 104, 127)


def _add_legend(slide: Any, chart_spec: dict[str, Any], chart_top: int, chart_height: int, region: tuple[int, int, int, int] | None = None) -> None:
    y_columns = chart_spec.get("y_columns") or []
    if not y_columns:
        return
    if region is None:
        legend_top = chart_top + chart_height + int(Inches(0.2))
        region = (int(Inches(0.6)), legend_top, int(Inches(8.2)), int(Inches(0.8)))
    _add_caption_box(slide, region[0], region[1], region[2], region[3], f"Legend: {', '.join(y_columns)}")


def _add_asset_result_content(
    slide: Any,
    slide_width: int,
    slide_height: int,
    chart_region: tuple[int, int, int, int],
    illustration_region: tuple[int, int, int, int],
    chart_path: str | Path | None,
    illustration_path: str | Path | None,
    title: str,
    summary: str,
    intent: dict[str, Any],
    title_region: tuple[int, int, int, int] | None = None,
    legend_region: tuple[int, int, int, int] | None = None,
) -> None:
    chart_type = intent.get("chart_type", "bar")
    model = intent.get("image_model", "local")
    style = intent.get("illustration_style", "auto")
    score = intent.get("clip_score")

    _add_title_block(slide, title, summary, title_region, slide_width, slide_height)

    if chart_path and _is_raster_image(chart_path) and Path(chart_path).exists():
        slide.shapes.add_picture(str(chart_path), Emu(chart_region[0]), Emu(chart_region[1]), width=Emu(chart_region[2]), height=Emu(chart_region[3]))
    else:
        _add_caption_box(slide, chart_region[0], chart_region[1], chart_region[2], int(slide_height * 0.12), "Chart preview is not available as a raster image yet.")

    legend_text = f"Chart type: {chart_type} | Source: {intent.get('source', 'local')} | Mode: {intent.get('semantic_mode', 'local')}"
    if legend_region is not None:
        _add_caption_box(slide, legend_region[0], legend_region[1], legend_region[2], legend_region[3], legend_text)
    else:
        _add_caption_box(slide, chart_region[0], chart_region[1] + chart_region[3] + int(slide_height * 0.02), chart_region[2], int(slide_height * 0.08), legend_text)

    if illustration_path and _is_raster_image(illustration_path) and Path(illustration_path).exists():
        slide.shapes.add_picture(str(illustration_path), Emu(illustration_region[0]), Emu(illustration_region[1]), width=Emu(illustration_region[2]), height=Emu(illustration_region[3]))
    else:
        _add_textbox(slide, illustration_region[0], illustration_region[1], illustration_region[2], int(slide_height * 0.08), "Illustration Area", 20, bold=True, color=(52, 82, 126))
        _add_caption_box(slide, illustration_region[0], illustration_region[1] + int(slide_height * 0.08), illustration_region[2], int(slide_height * 0.18), f"Style: {style} | Model: {model}\nMatch score: {score if score is not None else 'N/A'}\nTheme: {intent.get('visual_theme', 'Generated illustration preview')}")


def insert_chart_to_pptx(ppt_path: str | Path, chart_image_path: str | Path, slide_number: int, chart_title: str, chart_spec: dict[str, Any] | None = None, shapes: list[dict[str, Any]] | None = None, output_path: str | Path | None = None) -> InsertResult:
    _ensure_dependencies()
    source_path = Path(ppt_path)
    image_path = Path(chart_image_path)
    if not source_path.exists():
        raise FileNotFoundError(f"PPT file not found: {source_path}")
    if not image_path.exists():
        raise FileNotFoundError(f"Chart image not found: {image_path}")

    presentation = Presentation(str(source_path))
    if slide_number < 1 or slide_number > len(presentation.slides):
        raise ValueError(f"Slide number {slide_number} is out of range. This PPT has {len(presentation.slides)} slides.")

    slide = presentation.slides[slide_number - 1]
    slide_height = int(presentation.slide_height)
    anchor = _best_anchor(shapes or [])
    title_shape = _find_adjacent_text_shape(shapes or [], anchor, "above")
    legend_shape = _find_adjacent_text_shape(shapes or [], anchor, "below")
    removal_indexes = [anchor["index"]] + [shape["index"] for shape in [title_shape, legend_shape] if shape]
    replaced_table = anchor["index"] > 0
    _remove_shape_indexes(slide, removal_indexes)

    base_region = (
        anchor["left"],
        max(anchor["top"], int(Inches(1.1))),
        min(anchor["width"] or int(Inches(8.0)), int(Inches(8.6))),
        min(anchor["height"] or int(Inches(4.5)), int(Inches(4.8))),
    )
    title_region = _shape_bounds(title_shape) if title_shape else None
    legend_region = _shape_bounds(legend_shape) if legend_shape else None
    chart_left, chart_top, chart_width, chart_height = _compute_chart_region(base_region, slide_height, title_region, legend_region)

    _add_title(slide, chart_title, title_region)
    slide.shapes.add_picture(str(image_path), chart_left, chart_top, width=chart_width, height=chart_height)
    _add_legend(slide, chart_spec or {}, chart_top, chart_height, legend_region)

    destination = _derive_output_path(source_path, output_path)
    destination.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(str(destination))
    return InsertResult(str(destination), slide_number, int(chart_left), int(chart_top), int(chart_width), int(chart_height), replaced_table)


def insert_generated_assets(
    ppt_path: str | Path,
    output_path: str | Path,
    slide_number: int,
    chart_path: str | Path | None = None,
    illustration_path: str | Path | None = None,
    title: str = "",
    subtitle: str = "",
    intent: dict[str, Any] | None = None,
    shapes: list[dict[str, Any]] | None = None,
    layout_override: dict[str, Any] | None = None,
) -> Path:
    _ensure_dependencies()
    source = Path(ppt_path)
    target = Path(output_path)
    if not source.exists():
        raise FileNotFoundError(f"PPT file not found: {source}")

    presentation = Presentation(str(source))
    if slide_number < 1 or slide_number > len(presentation.slides):
        raise ValueError(f"Slide number {slide_number} is out of range.")

    slide = presentation.slides[slide_number - 1]
    slide_width = int(presentation.slide_width)
    slide_height = int(presentation.slide_height)
    intent = intent or {}

    anchor = _best_anchor(shapes or [])
    title_shape = _find_adjacent_text_shape(shapes or [], anchor, "above")
    legend_shape = _find_adjacent_text_shape(shapes or [], anchor, "below")
    title_region = _shape_bounds(title_shape) if title_shape else None
    legend_region = _shape_bounds(legend_shape) if legend_shape else None

    if shapes and anchor["index"]:
        base_region = (
            anchor["left"],
            max(anchor["top"], int(Inches(1.1))),
            min(anchor["width"] or int(Inches(8.0)), int(Inches(8.6))),
            min(anchor["height"] or int(Inches(4.5)), int(Inches(4.8))),
        )
        chart_anchor = _compute_chart_region(base_region, slide_height, title_region, legend_region)
    else:
        chart_anchor = _replace_table_region(slide)
        if chart_anchor is None:
            chart_anchor = _default_chart_region(slide_width, slide_height)
    manual_regions = _manual_override_regions(slide_width, slide_height, layout_override)
    if manual_regions:
        chart_region, illu_region = manual_regions
    else:
        chart_region, illu_region = _choose_asset_regions(slide_width, slide_height, shapes or [], chart_anchor, anchor["index"] if anchor["index"] else None)
    summary = subtitle or "Auto-generated chart and illustration preview"
    layout = _layout_metadata(chart_region, illu_region, shapes or [], slide_width, slide_height, anchor["index"] if anchor["index"] else None)
    if manual_regions:
        layout.update(
            {
                "insertion_mode": "manual_override",
                "manual_override": True,
                "original_slide_preserved": False,
                "layout_override": layout_override,
            }
        )
        _replace_table_shape(slide, anchor["index"])
        _add_asset_result_content(
            slide,
            slide_width,
            slide_height,
            chart_region,
            illu_region,
            chart_path,
            illustration_path,
            title or f"Slide {slide_number} chart result",
            summary,
            intent,
        )
    elif layout["layout_warning"]:
        result_slide = _add_blank_slide(presentation)
        result_chart_region = _result_slide_chart_region(slide_width, slide_height)
        result_illustration_region = _result_slide_illustration_region(slide_width, slide_height)
        layout.update(
            {
                "insertion_mode": "appendix",
                "original_slide_preserved": True,
                "result_slide_number": len(presentation.slides),
                "chart_region": {
                    "left": result_chart_region[0],
                    "top": result_chart_region[1],
                    "width": result_chart_region[2],
                    "height": result_chart_region[3],
                },
                "illustration_region": {
                    "left": result_illustration_region[0],
                    "top": result_illustration_region[1],
                    "width": result_illustration_region[2],
                    "height": result_illustration_region[3],
                },
                "pair_overlap": _overlap_area(result_chart_region, result_illustration_region),
            }
        )
        _add_asset_result_content(
            result_slide,
            slide_width,
            slide_height,
            result_chart_region,
            result_illustration_region,
            chart_path,
            illustration_path,
            title or f"Slide {slide_number} enhanced result",
            f"Original slide {slide_number} preserved. {summary}",
            intent,
        )
    else:
        removal_indexes = [anchor["index"]] + [shape["index"] for shape in [title_shape, legend_shape] if shape]
        _remove_shape_indexes(slide, removal_indexes)
        _add_asset_result_content(
            slide,
            slide_width,
            slide_height,
            chart_region,
            illu_region,
            chart_path,
            illustration_path,
            title or f"Slide {slide_number} chart result",
            summary,
            intent,
            title_region=title_region,
            legend_region=legend_region,
        )

    intent["layout"] = layout

    target.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(str(target))
    return target
