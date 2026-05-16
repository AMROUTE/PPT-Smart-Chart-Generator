from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path
import re
from typing import Any

from PIL import Image, ImageDraw, ImageFont

try:
    import pandas as pd
except ModuleNotFoundError:  # pragma: no cover
    pd = None

try:
    from pptx import Presentation
except ModuleNotFoundError:  # pragma: no cover
    Presentation = None


@dataclass
class ParsedSlideContent:
    slide_number: int
    text_content: str
    tables: list[dict[str, Any]]
    shapes: list[dict[str, Any]]

    def to_dict(self) -> dict[str, Any]:
        return {
            "slide_number": self.slide_number,
            "text_content": self.text_content,
            "tables": self.tables,
            "shapes": self.shapes,
        }


@dataclass
class SlidePreview:
    slide_number: int
    slide_count: int
    output_path: str

    def to_dict(self) -> dict[str, Any]:
        return {
            "slide_number": self.slide_number,
            "slide_count": self.slide_count,
            "output_path": self.output_path,
        }


@dataclass
class SlideOutlineItem:
    slide_number: int
    text_content: str
    table_count: int
    shape_count: int
    table_titles: list[str]

    def to_dict(self) -> dict[str, Any]:
        return {
            "slide_number": self.slide_number,
            "text_content": self.text_content,
            "table_count": self.table_count,
            "shape_count": self.shape_count,
            "table_titles": self.table_titles,
        }


def normalize_header(value: str | None, index: int) -> str:
    text = (value or "").strip()
    return text or f"column_{index + 1}"


def dataframe_to_records(dataframe: pd.DataFrame) -> list[list[Any]]:
    return dataframe.where(pd.notnull(dataframe), None).values.tolist()


def _normalized_cell_text(cell: Any) -> str:
    text = cell.text.strip()
    lines = [line.strip() for line in text.splitlines() if line.strip()]
    return lines[0] if lines else text


def _cell_merge_meta(cell: Any, row_index: int, col_index: int) -> dict[str, Any]:
    return {
        "row": row_index,
        "col": col_index,
        "text": _normalized_cell_text(cell),
        "is_merge_origin": bool(getattr(cell, "is_merge_origin", False)),
        "is_spanned": bool(getattr(cell, "is_spanned", False)),
        "row_span": int(getattr(cell, "span_height", 1) or 1),
        "col_span": int(getattr(cell, "span_width", 1) or 1),
    }


def extract_table_matrix(table: Any) -> list[list[dict[str, Any]]]:
    matrix: list[list[dict[str, Any]]] = []
    for row_index, row in enumerate(table.rows):
        matrix.append([_cell_merge_meta(cell, row_index, col_index) for col_index, cell in enumerate(row.cells)])
    return matrix


def _resolve_visible_value(matrix: list[list[dict[str, Any]]], row_index: int, col_index: int) -> str:
    cell = matrix[row_index][col_index]
    if not cell["is_spanned"]:
        return cell["text"]

    for source_row in range(row_index, -1, -1):
        for source_col in range(col_index, -1, -1):
            source = matrix[source_row][source_col]
            if not source["is_merge_origin"]:
                continue
            row_end = source_row + source["row_span"] - 1
            col_end = source_col + source["col_span"] - 1
            if source_row <= row_index <= row_end and source_col <= col_index <= col_end:
                return source["text"]
    return cell["text"]


def table_to_dataframe(table: Any) -> pd.DataFrame:
    if pd is None:
        raise ModuleNotFoundError("pandas is required for table_to_dataframe.")

    raw_matrix = extract_table_matrix(table)
    matrix = [
        [_resolve_visible_value(raw_matrix, row_index, col_index) for col_index in range(len(row))]
        for row_index, row in enumerate(raw_matrix)
    ]
    if not matrix:
        return pd.DataFrame()

    normalized_matrix: list[list[str]] = []
    for row_index, row in enumerate(matrix):
        filled_row: list[str] = []
        for col_index, value in enumerate(row):
            normalized = value
            if not normalized and col_index > 0:
                normalized = filled_row[col_index - 1]
            if not normalized and row_index > 0:
                normalized = normalized_matrix[row_index - 1][col_index]
            filled_row.append(normalized)
        normalized_matrix.append(filled_row)

    headers = [normalize_header(value, index) for index, value in enumerate(normalized_matrix[0])]
    deduped_headers: list[str] = []
    header_counts: dict[str, int] = {}
    for header in headers:
        count = header_counts.get(header, 0)
        deduped_headers.append(header if count == 0 else f"{header}_{count + 1}")
        header_counts[header] = count + 1

    return pd.DataFrame(normalized_matrix[1:], columns=deduped_headers)


def extract_text_from_shape(shape: Any) -> str:
    if not getattr(shape, "has_text_frame", False):
        return ""
    return shape.text.strip()


def describe_shape(shape: Any, index: int) -> dict[str, Any]:
    shape_type = getattr(getattr(shape, "shape_type", None), "name", str(getattr(shape, "shape_type", "")))
    payload = {
        "index": index,
        "name": getattr(shape, "name", f"shape_{index}"),
        "shape_type": shape_type,
        "has_text": bool(getattr(shape, "has_text_frame", False)),
        "has_table": bool(getattr(shape, "has_table", False)),
        "left": int(getattr(shape, "left", 0)),
        "top": int(getattr(shape, "top", 0)),
        "width": int(getattr(shape, "width", 0)),
        "height": int(getattr(shape, "height", 0)),
    }
    text = extract_text_from_shape(shape)
    if text:
        payload["text"] = text
    return payload


def extract_tables(slide: Any) -> list[dict[str, Any]]:
    tables: list[dict[str, Any]] = []
    table_index = 1
    for shape in slide.shapes:
        if not getattr(shape, "has_table", False):
            continue

        table_matrix = extract_table_matrix(shape.table)
        dataframe = table_to_dataframe(shape.table)
        raw_matrix = [[cell.text.strip() for cell in row.cells] for row in shape.table.rows]
        merged_hints: list[dict[str, Any]] = []
        for row_index, row in enumerate(raw_matrix):
            for col_index, value in enumerate(row):
                if value:
                    continue
                if col_index > 0 and row[col_index - 1]:
                    merged_hints.append({"row": row_index, "col": col_index, "merge_origin": [row_index, col_index - 1], "direction": "horizontal"})
                elif row_index > 0 and raw_matrix[row_index - 1][col_index]:
                    merged_hints.append({"row": row_index, "col": col_index, "merge_origin": [row_index - 1, col_index], "direction": "vertical"})

        tables.append(
            {
                "title": getattr(shape, "name", f"table_{table_index}"),
                "columns": dataframe.columns.tolist(),
                "rows": dataframe_to_records(dataframe),
                "dataframe": dataframe,
                "cell_matrix": table_matrix,
                "merge_hints": merged_hints,
                "raw_matrix": raw_matrix,
            }
        )
        table_index += 1
    return tables


def extract_shapes(slide: Any) -> list[dict[str, Any]]:
    return [describe_shape(shape, index) for index, shape in enumerate(slide.shapes, start=1)]


def infer_table_from_text_blocks(text_blocks: list[str]) -> list[dict[str, Any]]:
    joined = "\n".join(text_blocks)
    matches = re.findall(r"([A-Za-z\u4e00-\u9fff][A-Za-z0-9\u4e00-\u9fff\s]{0,24})[:\s]*(-?\d+(?:\.\d+)?)", joined)
    if not matches:
        matches = re.findall(r"((?:20\d{2}|Q[1-4]|[1-9]|1[0-2])\s*[,:-]?\s*)(-?\d+(?:\.\d+)?)", joined)
        matches = [(label, value) for label, value in matches]
    if not matches:
        return []
    rows = [[label.strip(), float(value)] for label, value in matches]
    return [{"title": "text_inferred_table", "columns": ["category", "value"], "rows": rows, "merge_hints": [], "raw_matrix": [["category", "value"], *rows], "cell_matrix": []}]


def _ensure_presentation_dependency() -> None:
    if Presentation is None:
        raise ModuleNotFoundError("python-pptx is required for PPT parsing.")


def _load_presentation(ppt_path: str | Path) -> Any:
    _ensure_presentation_dependency()
    presentation_path = Path(ppt_path)
    if not presentation_path.exists():
        raise FileNotFoundError(f"PPT file not found: {presentation_path}")
    return Presentation(str(presentation_path))


def _validate_slide_number(presentation: Any, slide_number: int) -> None:
    if slide_number < 1 or slide_number > len(presentation.slides):
        raise ValueError(f"Slide number {slide_number} is out of range. This PPT has {len(presentation.slides)} slides.")


def _build_parsed_slide(slide: Any, slide_number: int) -> ParsedSlideContent:
    tables = extract_tables(slide)
    shapes = extract_shapes(slide)
    text_blocks = [shape["text"] for shape in shapes if shape.get("text")]
    if not tables:
        tables = infer_table_from_text_blocks(text_blocks)
    return ParsedSlideContent(slide_number=slide_number, text_content="\n".join(text_blocks), tables=tables, shapes=shapes)


def extract_slide_content_from_presentation(presentation: Any, slide_number: int) -> ParsedSlideContent:
    _validate_slide_number(presentation, slide_number)
    return _build_parsed_slide(presentation.slides[slide_number - 1], slide_number)


def extract_multiple_slide_contents(ppt_path: str | Path, slide_numbers: list[int] | None = None) -> dict[int, ParsedSlideContent]:
    presentation = _load_presentation(ppt_path)
    total_slides = len(presentation.slides)
    requested = slide_numbers or list(range(1, total_slides + 1))
    parsed: dict[int, ParsedSlideContent] = {}
    for slide_number in requested:
        _validate_slide_number(presentation, slide_number)
        parsed[slide_number] = _build_parsed_slide(presentation.slides[slide_number - 1], slide_number)
    return parsed


def extract_slide_content(ppt_path: str | Path, slide_number: int) -> ParsedSlideContent:
    presentation = _load_presentation(ppt_path)
    return extract_slide_content_from_presentation(presentation, slide_number)


def get_slide_count(ppt_path: str | Path) -> int:
    presentation = _load_presentation(ppt_path)
    return len(presentation.slides)


def render_slide_preview(ppt_path: str | Path, slide_number: int, output_path: str | Path) -> SlidePreview:
    presentation = _load_presentation(ppt_path)
    _validate_slide_number(presentation, slide_number)

    slide_count = len(presentation.slides)
    slide = presentation.slides[slide_number - 1]
    slide_width = int(presentation.slide_width) or 1
    slide_height = int(presentation.slide_height) or 1
    canvas_width = 1280
    canvas_height = max(720, int(canvas_width * slide_height / slide_width))
    scale_x = canvas_width / slide_width
    scale_y = canvas_height / slide_height

    image = Image.new("RGB", (canvas_width, canvas_height), "#f8fbff")
    draw = ImageDraw.Draw(image)
    title_font = ImageFont.load_default()
    body_font = ImageFont.load_default()

    draw.rounded_rectangle((20, 20, canvas_width - 20, canvas_height - 20), radius=28, outline="#c6d7ea", width=3)

    for shape in slide.shapes:
        left = int(getattr(shape, "left", 0) * scale_x)
        top = int(getattr(shape, "top", 0) * scale_y)
        width = max(18, int(getattr(shape, "width", 0) * scale_x))
        height = max(18, int(getattr(shape, "height", 0) * scale_y))
        box = (left, top, left + width, top + height)
        if getattr(shape, "has_table", False):
            draw.rounded_rectangle(box, radius=18, fill="#ffffff", outline="#3b82f6", width=3)
            rows = len(shape.table.rows)
            cols = len(shape.table.columns)
            if rows > 0:
                row_height = height / rows
                for row_index in range(1, rows):
                    y = top + int(row_height * row_index)
                    draw.line((left, y, left + width, y), fill="#cfe1f5", width=2)
            if cols > 0:
                col_width = width / cols
                for col_index in range(1, cols):
                    x = left + int(col_width * col_index)
                    draw.line((x, top, x, top + height), fill="#cfe1f5", width=2)
            header = shape.table.cell(0, 0).text.strip() if rows and cols else "Table"
            draw.text((left + 12, top + 10), header[:18] or "Table", fill="#214872", font=body_font)
            continue
        if getattr(shape, "has_text_frame", False):
            text = shape.text.strip()
            draw.rounded_rectangle(box, radius=18, fill="#ffffff", outline="#d5dfeb", width=2)
            lines = [line.strip() for line in text.splitlines() if line.strip()]
            if not lines:
                continue
            draw.text((left + 12, top + 10), lines[0][:40], fill="#102033", font=title_font)
            for line_index, line in enumerate(lines[1:4], start=1):
                draw.text((left + 12, top + 10 + line_index * 18), line[:42], fill="#4b627a", font=body_font)
            continue
        draw.rounded_rectangle(box, radius=18, fill="#eef5ff", outline="#d5dfeb", width=2)

    output = Path(output_path)
    output.parent.mkdir(parents=True, exist_ok=True)
    image.save(output)
    return SlidePreview(slide_number=slide_number, slide_count=slide_count, output_path=str(output))


def parse_presentation_outline(ppt_path: str | Path) -> list[SlideOutlineItem]:
    presentation = _load_presentation(ppt_path)
    items: list[SlideOutlineItem] = []
    for slide_index in range(1, len(presentation.slides) + 1):
        parsed = extract_slide_content_from_presentation(presentation, slide_index)
        items.append(
            SlideOutlineItem(
                slide_number=slide_index,
                text_content=parsed.text_content[:600],
                table_count=len(parsed.tables),
                shape_count=len(parsed.shapes),
                table_titles=[table.get("title", f"table_{table_index + 1}") for table_index, table in enumerate(parsed.tables)],
            )
        )
    return items
