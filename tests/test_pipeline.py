import os
import tempfile
import unittest
from pathlib import Path
from types import SimpleNamespace
from unittest.mock import Mock, patch

os.environ["ENABLE_QWEN_API"] = "0"
os.environ["DATABASE_PATH"] = str(Path(tempfile.gettempdir()) / "codex-test-app.db")

try:
    from PIL import Image
    from pptx import Presentation
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
    from pptx.util import Inches
except ModuleNotFoundError:
    Image = None
    Presentation = None
    MSO_AUTO_SHAPE_TYPE = None
    Inches = None

from fastapi.testclient import TestClient

from backend.app import create_app
from backend.chart_generator import generate_chart
from backend.database import authenticate_or_create_user, fetch_processing_job, init_db
from backend.image_clients import _resolve_flux_result_url, generate_flux_image, generate_wanx_image
from backend.insert_to_pptx import _choose_asset_regions, _overlap_area, insert_chart_to_pptx, insert_generated_assets
from backend.pipeline import PIPELINE_NODES, _infer_illustration_context, _recommend_chart_intent, _select_illustration_composition_variant, export_pipeline_mermaid, run_pipeline
from backend.ppt_parser import extract_slide_content, table_to_dataframe
from backend.schemas import PipelineInput
from backend.services import (
    allowed_file,
    build_file_metadata,
    build_health_payload,
    build_slide_preview,
    extract_records_from_text,
    normalize_chart_theme,
    normalize_chart_type_override,
    normalize_image_model,
    normalize_illustration_style,
    parse_presentation_slides,
    apply_batch_layout_overrides,
    path_to_asset_url,
    process_demo_text,
    process_ppt_batch,
    process_local_ppt,
    process_local_ppt_batch,
)


class PipelineTests(unittest.TestCase):
    def test_pipeline_nodes_are_defined_for_week_two(self):
        self.assertEqual(PIPELINE_NODES, ["parse_ppt", "semantic_analysis", "generate_chart", "generate_illustration", "save_pptx"])

    def test_mermaid_definition_contains_full_flow(self):
        mermaid = export_pipeline_mermaid()
        self.assertIn("parse_ppt", mermaid)
        self.assertIn("save_pptx", mermaid)

    def test_run_pipeline_returns_expected_placeholders(self):
        result = run_pipeline(PipelineInput(ppt_path="demo.pptx", current_slide=2, request_id="test-run"))
        self.assertIn(result["intent"]["chart_type"], {"bar", "line", "pie", "scatter", "heatmap"})
        self.assertIn("chart_slide_2", result["chart_image"])
        self.assertTrue(result["final_pptx_path"].endswith("demo_enhanced.pptx"))
        self.assertGreaterEqual(len(result["logs"]), 6)
        self.assertEqual(result["status"], "completed")
        self.assertEqual(result["progress"], 100)
        self.assertEqual(len(result["stage_history"]), 5)


class ServiceTests(unittest.TestCase):
    @classmethod
    def setUpClass(cls):
        init_db()

    def test_allowed_file_only_accepts_pptx(self):
        self.assertTrue(allowed_file("demo.pptx"))
        self.assertFalse(allowed_file("demo.pdf"))

    def test_file_metadata_tracks_slide_number(self):
        with tempfile.NamedTemporaryFile(suffix=".pptx") as tmp:
            metadata = build_file_metadata(Path(tmp.name), 3)
            self.assertEqual(metadata["slide_number"], 3)
            self.assertEqual(metadata["suffix"], ".pptx")

    def test_process_local_ppt_runs_pipeline(self):
        if Presentation is None:
            self.skipTest("python-pptx is not installed")
        tmp_path = Path(tempfile.gettempdir()) / "codex-test-source.pptx"
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.add_textbox(1000000, 1000000, 4000000, 600000).text_frame.text = "Revenue trend analysis"
        table = slide.shapes.add_table(4, 2, 1000000, 1800000, 4000000, 2000000).table
        table.cell(0, 0).text = "Quarter"
        table.cell(0, 1).text = "Revenue"
        table.cell(1, 0).text = "Q1"
        table.cell(1, 1).text = "120"
        table.cell(2, 0).text = "Q2"
        table.cell(2, 1).text = "150"
        table.cell(3, 0).text = "Q3"
        table.cell(3, 1).text = "180"
        prs.save(tmp_path)
        try:
            payload = process_local_ppt(tmp_path, 1, chart_type_override="line", chart_theme="business", illustration_style="tech", image_model="flux")
            self.assertEqual(payload["file"]["slide_number"], 1)
            self.assertIn("pipeline", payload)
            self.assertIn("chart_image_url", payload["pipeline"])
            self.assertTrue(Path(payload["pipeline"]["final_pptx_path"]).exists())
            self.assertEqual(payload["pipeline"]["chart_spec"]["theme"], "business")
            self.assertIn("layout", payload["pipeline"]["intent"])
            self.assertIn("layout_warning", payload["pipeline"]["intent"]["layout"])
        finally:
            tmp_path.unlink(missing_ok=True)

    def test_process_local_ppt_batch_reuses_single_batch_output(self):
        if Presentation is None:
            self.skipTest("python-pptx is not installed")
        tmp_path = Path(tempfile.gettempdir()) / "codex-test-batch-source.pptx"
        prs = Presentation()
        for title, values in [("Revenue trend", [120, 150, 180]), ("Market share", [35, 25, 40])]:
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            slide.shapes.add_textbox(1000000, 1000000, 4000000, 600000).text_frame.text = title
            table = slide.shapes.add_table(4, 2, 1000000, 1800000, 4000000, 2000000).table
            table.cell(0, 0).text = "Label"
            table.cell(0, 1).text = "Value"
            table.cell(1, 0).text = "A"
            table.cell(1, 1).text = str(values[0])
            table.cell(2, 0).text = "B"
            table.cell(2, 1).text = str(values[1])
            table.cell(3, 0).text = "C"
            table.cell(3, 1).text = str(values[2])
        prs.save(tmp_path)
        try:
            payload = process_local_ppt_batch(tmp_path, [1, 2], chart_theme="academic", illustration_style="tech")
            self.assertEqual(payload["processed_count"], 2)
            self.assertEqual(payload["slide_numbers"], [1, 2])
            self.assertTrue(Path(payload["final_pptx_path"]).exists())
            self.assertTrue(all(slide["final_pptx_path"] == payload["final_pptx_path"] for slide in payload["slides"]))
            self.assertTrue(all(slide["chart_spec"]["theme"] == "academic" for slide in payload["slides"]))
        finally:
            tmp_path.unlink(missing_ok=True)

    def test_process_local_ppt_batch_reports_partial_failures(self):
        with patch("backend.ppt_parser.extract_multiple_slide_contents") as mock_extract, patch("backend.ppt_parser.get_slide_count", return_value=2), patch("backend.pipeline.run_pipeline") as mock_run:
            from backend.services import process_local_ppt_batch

            tmp_path = Path(tempfile.gettempdir()) / "codex-test-local-batch-failure.pptx"
            tmp_path.write_bytes(b"ppt")
            parsed = SimpleNamespace(text_content="Revenue trend", tables=[{"title": "t", "columns": ["label", "value"], "rows": [["A", 1]], "cell_matrix": [], "merge_hints": [], "raw_matrix": []}], shapes=[])
            mock_extract.return_value = {1: parsed, 2: parsed}
            mock_run.side_effect = [
                {"request_id": "batch-x-s1", "current_slide": 1, "status": "completed", "final_pptx_path": str(tmp_path), "chart_spec": {"theme": "tech"}, "intent": {}, "chart_image": "", "illustration_image": ""},
                RuntimeError("slide 2 failed"),
            ]
            try:
                payload = process_local_ppt_batch(tmp_path, [1, 2])
                self.assertEqual(payload["batch"]["status"], "partial")
                self.assertEqual(payload["batch"]["success_count"], 1)
                self.assertEqual(payload["batch"]["failure_count"], 1)
                self.assertEqual(payload["result_summary"]["failure_count"], 1)
                self.assertEqual(payload["batch"]["slides"][0]["result_level"], "pass")
                self.assertEqual(payload["batch"]["slides"][1]["status"], "failed")
                self.assertEqual(payload["batch"]["slides"][1]["result_level"], "failed")
                self.assertIn("slide 2 failed", payload["batch"]["slides"][1]["error"])
            finally:
                tmp_path.unlink(missing_ok=True)

    def test_process_ppt_batch_runs_selected_slide_range(self):
        if Presentation is None:
            self.skipTest("python-pptx is not installed")
        tmp_path = Path(tempfile.gettempdir()) / "codex-test-batch-source.pptx"
        prs = Presentation()
        for slide_index in range(1, 4):
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            slide.shapes.add_textbox(1000000, 700000, 5000000, 700000).text_frame.text = f"第 {slide_index} 页营收趋势"
            table = slide.shapes.add_table(4, 2, 1000000, 1800000, 4000000, 2000000).table
            table.cell(0, 0).text = "季度"
            table.cell(0, 1).text = "营收"
            table.cell(1, 0).text = "Q1"
            table.cell(1, 1).text = str(100 + slide_index)
            table.cell(2, 0).text = "Q2"
            table.cell(2, 1).text = str(140 + slide_index)
            table.cell(3, 0).text = "Q3"
            table.cell(3, 1).text = str(180 + slide_index)
        prs.save(tmp_path)
        try:
            payload = process_ppt_batch(
                tmp_path,
                slide_start=1,
                slide_end=2,
                chart_type_override="line",
                illustration_style="tech",
                image_model="local",
            )
            self.assertEqual(payload["batch"]["total_slides"], 2)
            self.assertEqual(payload["batch"]["success_count"], 2)
            self.assertEqual(payload["batch"]["failure_count"], 0)
            self.assertEqual(payload["result_summary"]["failure_count"], 0)
            self.assertEqual([item["slide_number"] for item in payload["batch"]["slides"]], [1, 2])
            self.assertTrue(all("request_id" in item for item in payload["batch"]["slides"]))
            self.assertTrue(all(item["result_level"] in {"pass", "warning"} for item in payload["batch"]["slides"]))
            self.assertTrue(Path(payload["batch"]["final_pptx_path"]).exists())
            self.assertTrue(payload["batch"]["final_pptx_url"].startswith("/assets/outputs/"))
            self.assertTrue(payload["batch"]["slides"][0]["pipeline"]["chart_image_url"].startswith("/assets/outputs/"))
        finally:
            tmp_path.unlink(missing_ok=True)

    def test_process_ppt_batch_skips_empty_slides(self):
        if Presentation is None:
            self.skipTest("python-pptx is not installed")
        tmp_path = Path(tempfile.gettempdir()) / "codex-test-batch-empty-slide.pptx"
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.add_textbox(1000000, 700000, 5000000, 700000).text_frame.text = "Revenue trend"
        table = slide.shapes.add_table(4, 2, 1000000, 1800000, 4000000, 2000000).table
        table.cell(0, 0).text = "Quarter"
        table.cell(0, 1).text = "Revenue"
        table.cell(1, 0).text = "Q1"
        table.cell(1, 1).text = "100"
        table.cell(2, 0).text = "Q2"
        table.cell(2, 1).text = "140"
        table.cell(3, 0).text = "Q3"
        table.cell(3, 1).text = "180"
        prs.slides.add_slide(prs.slide_layouts[6])
        prs.save(tmp_path)
        try:
            payload = process_ppt_batch(tmp_path, slide_start=1, slide_end=2, image_model="local")
            self.assertEqual(payload["batch"]["total_slides"], 2)
            self.assertEqual(payload["batch"]["success_count"], 1)
            self.assertEqual(payload["batch"]["skipped_count"], 1)
            self.assertEqual(payload["batch"]["slides"][1]["status"], "skipped")
            self.assertEqual(payload["batch"]["slides"][1]["reason"], "Empty slide skipped.")
            self.assertTrue(Path(payload["batch"]["final_pptx_path"]).exists())
        finally:
            tmp_path.unlink(missing_ok=True)

    def test_process_ppt_batch_uses_shapes_for_layout_writeback(self):
        if Presentation is None or Image is None:
            self.skipTest("python-pptx and Pillow are required")
        tmp_path = Path(tempfile.gettempdir()) / "codex-test-batch-dense-layout.pptx"
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        for row in range(4):
            for col in range(3):
                box = slide.shapes.add_textbox(
                    Inches(0.4 + col * 3.0),
                    Inches(0.9 + row * 1.25),
                    Inches(2.45),
                    Inches(0.85),
                )
                box.text_frame.text = f"Metric {row}-{col}: {100 + row * 20 + col * 5}"
        prs.save(tmp_path)
        try:
            payload = process_ppt_batch(tmp_path, slide_start=1, slide_end=1, image_model="local")
            final_path = Path(payload["batch"]["final_pptx_path"])
            self.assertTrue(final_path.exists())
            enhanced = Presentation(str(final_path))
            self.assertEqual(len(enhanced.slides), 2)
            layout = payload["batch"]["slides"][0]["pipeline"]["intent"]["layout"]
            self.assertEqual(layout["insertion_mode"], "appendix")
            self.assertTrue(layout["original_slide_preserved"])
        finally:
            tmp_path.unlink(missing_ok=True)

    def test_build_slide_preview_returns_preview_asset(self):
        if Presentation is None:
            self.skipTest("python-pptx is not installed")
        tmp_path = Path(tempfile.gettempdir()) / "codex-test-preview.pptx"
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.add_textbox(1000000, 1000000, 5000000, 1200000).text_frame.text = "First slide preview test"
        prs.save(tmp_path)
        try:
            payload = build_slide_preview(1, file_path=tmp_path)
            self.assertEqual(payload["slide_number"], 1)
            self.assertEqual(payload["slide_count"], 1)
            self.assertTrue(Path(payload["preview_image"]).exists())
            self.assertTrue(payload["preview_image_url"].startswith("/assets/outputs/"))
        finally:
            tmp_path.unlink(missing_ok=True)

    def test_parse_presentation_slides_returns_outline(self):
        if Presentation is None:
            self.skipTest("python-pptx is not installed")
        tmp_path = Path(tempfile.gettempdir()) / "codex-test-outline.pptx"
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[6]).shapes.add_textbox(1000000, 1000000, 5000000, 1200000).text_frame.text = "Slide one"
        prs.slides.add_slide(prs.slide_layouts[6]).shapes.add_textbox(1000000, 1000000, 5000000, 1200000).text_frame.text = "Slide two"
        prs.save(tmp_path)
        try:
            payload = parse_presentation_slides(file_path=tmp_path)
            self.assertEqual(payload["slide_count"], 2)
            self.assertEqual(len(payload["slides"]), 2)
            self.assertEqual(payload["slides"][0]["slide_number"], 1)
            self.assertIn("diagnostics", payload["slides"][0])
            self.assertIn("picture_count", payload["slides"][0])
            self.assertFalse(payload["slides"][0]["is_empty"])
        finally:
            tmp_path.unlink(missing_ok=True)

    def test_extract_records_from_text_supports_demo_mode(self):
        records = extract_records_from_text("Revenue: 120\nProfit: 45")
        self.assertEqual(len(records), 2)
        self.assertEqual(records[0]["category"], "Revenue")

    def test_extract_records_from_text_preserves_numeric_labels(self):
        records = extract_records_from_text("2020: 100\n2021: 150\n2022: 220\n2023: 300")
        self.assertEqual(
            records,
            [
                {"category": "2020", "value": 100.0},
                {"category": "2021", "value": 150.0},
                {"category": "2022", "value": 220.0},
                {"category": "2023", "value": 300.0},
            ],
        )

    def test_extract_records_from_text_preserves_range_labels(self):
        records = extract_records_from_text("0-10分: 2\n10-20分: 4\n20-30分: 9")
        self.assertEqual([record["category"] for record in records], ["0-10分", "10-20分", "20-30分"])
        self.assertEqual([record["value"] for record in records], [2.0, 4.0, 9.0])

    def test_extract_records_from_text_pairs_repeated_metrics_for_correlation(self):
        records = extract_records_from_text("广告投入: 10\n销售额: 80\n广告投入: 20\n销售额: 120\n广告投入: 30\n销售额: 170")
        self.assertEqual(records[0], {"point": "Point 1", "广告投入": 10.0, "销售额": 80.0})
        self.assertEqual(records[2], {"point": "Point 3", "广告投入": 30.0, "销售额": 170.0})

    def test_process_demo_text_uses_real_two_axis_scatter_data(self):
        payload = process_demo_text("广告投入: 10\n销售额: 80\n广告投入: 20\n销售额: 120\n广告投入: 30\n销售额: 170\n广告投入越高销售额越高", image_model="local")
        spec = payload["pipeline"]["chart_spec"]
        self.assertEqual(payload["pipeline"]["intent"]["intent_category"], "correlation")
        self.assertEqual(spec["chart_type"], "scatter")
        self.assertEqual(spec["y_columns"], ["广告投入", "销售额"])
        self.assertFalse(any("synthetic" in warning.lower() for warning in spec["warnings"]))
        self.assertIn("scatter_real_xy", spec["render_notes"])
        self.assertIn("scatter_trendline", spec["render_notes"])

    def test_process_demo_text_returns_preview_assets(self):
        payload = process_demo_text("Q1: 12\nQ2: 18\nQ3: 26", chart_type_override="pie", chart_theme="minimal", illustration_style="business", image_model="wanx")
        self.assertEqual(payload["pipeline"]["status"], "completed")
        self.assertTrue(payload["pipeline"]["chart_image_url"].startswith("/assets/outputs/"))
        self.assertEqual(payload["pipeline"]["intent"]["chart_type"], "pie")
        self.assertIn("intent_category", payload["pipeline"]["intent"])
        self.assertIn("recommendation_confidence", payload["pipeline"]["intent"])
        self.assertIn("recommendation_signals", payload["pipeline"]["intent"])
        self.assertIn("clip_score", payload["pipeline"]["illustration_meta"])
        self.assertIn("quality_score", payload["pipeline"]["chart_spec"])
        self.assertGreater(payload["pipeline"]["chart_spec"]["quality_score"], 0)
        self.assertIn("quality_checks", payload["pipeline"]["chart_spec"])
        self.assertEqual(payload["pipeline"]["chart_spec"]["theme"], "minimal")
        self.assertIn(payload["pipeline"]["illustration_meta"]["generation_source"], {"local", "wanx", "flux"})

    def test_process_demo_text_falls_back_when_remote_image_model_is_unavailable(self):
        payload = process_demo_text("Revenue: 120\nCost: 80\nProfit: 40", illustration_style="tech", image_model="wanx")
        self.assertEqual(payload["pipeline"]["status"], "completed")
        meta = payload["pipeline"]["illustration_meta"]
        self.assertEqual(meta["generation_source"], "local")
        self.assertEqual(meta["requested_image_model"], "wanx")
        self.assertTrue(meta["external_provider_requested"])
        self.assertEqual(meta["external_provider"], "wanx")
        self.assertEqual(meta["resolved_image_source"], "local")
        self.assertTrue(meta["fallback_to_local"])
        self.assertTrue(meta["generation_warning"])
        self.assertIn("tech_device_cloud", meta["local_render_features"])
        self.assertTrue(Path(payload["pipeline"]["illustration_image"]).exists())

    def test_process_demo_text_initializes_empty_database(self):
        database_path = Path(os.environ["DATABASE_PATH"])
        database_path.unlink(missing_ok=True)
        payload = process_demo_text("Cold start: 12\nReady: 18", image_model="local")
        request_id = payload["pipeline"]["request_id"]
        job = fetch_processing_job(request_id)
        self.assertIsNotNone(job)
        self.assertEqual(job["request_id"], request_id)
        self.assertEqual(job["source_type"], "demo")

    def test_process_demo_text_regenerates_low_score_local_illustration(self):
        payload = process_demo_text("Alpha\nBeta", illustration_style="auto", image_model="local")
        meta = payload["pipeline"]["illustration_meta"]
        self.assertTrue(meta["regenerated"])
        self.assertEqual(meta["regenerate_attempts"], 1)
        self.assertEqual(meta["regenerate_action"], "local_refined_prompt")
        self.assertLess(meta["initial_clip_score"], meta["score_threshold"])
        self.assertGreaterEqual(meta["clip_score"], meta["score_threshold"])
        self.assertFalse(meta["regenerate_hint"])
        self.assertIn("quality_components", meta)
        self.assertTrue(meta["quality_components"]["negative_prompt"])
        self.assertIn("infographic", meta["negative_prompt_terms"])
        self.assertTrue(any(feature.startswith("business_") for feature in meta["local_render_features"]))
        self.assertIn("negative prompt", " ".join(meta["prompt_quality_notes"]).lower())
        self.assertIn("illustration_prompt_retry", payload["pipeline"])
        self.assertIn("Negative prompt: no charts", payload["pipeline"]["illustration_prompt_retry"])

    def test_process_demo_text_records_local_illustration_style_features(self):
        payload = process_demo_text("Healthcare visits: 42\nFollow-up calls: 18", illustration_style="medical", image_model="local")
        meta = payload["pipeline"]["illustration_meta"]
        self.assertEqual(meta["generation_source"], "local")
        self.assertIn("local_scene_preview", meta["local_render_features"])
        self.assertIn("human_subjects", meta["local_render_features"])
        self.assertIn("medical_care_symbol", meta["local_render_features"])
        self.assertIn(meta["composition_variant"], {"duo_panel", "full_scene", "spotlight", "diagonal_workshop"})
        self.assertTrue(any(feature.startswith("layout_variant_") for feature in meta["local_render_features"]))
        self.assertTrue(Path(payload["pipeline"]["illustration_image"]).exists())

    def test_illustration_composition_variant_is_stable_and_content_sensitive(self):
        first = _select_illustration_composition_variant("medical visits", "medical", "local", "slide-1")
        second = _select_illustration_composition_variant("medical visits", "medical", "local", "slide-1")
        variants = {
            _select_illustration_composition_variant(f"theme-{index}", "business", "local", f"slide-{index}")
            for index in range(8)
        }
        self.assertEqual(first, second)
        self.assertGreaterEqual(len(variants), 2)

    def test_illustration_context_avoids_generic_office_for_business_topics(self):
        samples = [
            (
                "华东: 35\n华南: 25\n华北: 20\n西南: 20\n请展示各区域市场份额占比",
                "regional",
                "business_regional_network",
            ),
            (
                "产品A: 120\n产品B: 95\n产品C: 150\n产品D: 110\n比较四个产品销量差异",
                "product",
                "business_product_showroom",
            ),
            (
                "广告投入: 10\n销售额: 80\n广告投入: 20\n销售额: 120\n广告投入与销售额呈正相关",
                "campaign",
                "business_marketing_studio",
            ),
        ]
        for text, theme_marker, expected_feature in samples:
            records = extract_records_from_text(text)
            columns = list(records[0].keys())
            rows = [[record[column] for column in columns] for record in records]
            recommendation = _recommend_chart_intent(text, columns, rows)
            context = _infer_illustration_context(text, columns, rows, recommendation, "auto")
            self.assertIn(theme_marker, context["visual_theme"].lower())
            self.assertNotIn("office collaboration", context["visual_theme"].lower())
            payload = process_demo_text(text, illustration_style="auto", image_model="local")
            meta = payload["pipeline"]["illustration_meta"]
            self.assertIn(expected_feature, meta["local_render_features"])

    def test_process_demo_text_persists_pipeline_metadata_for_logs(self):
        payload = process_demo_text("Alpha\nBeta", illustration_style="auto", image_model="local")
        request_id = payload["pipeline"]["request_id"]
        job = fetch_processing_job(request_id)
        self.assertIsNotNone(job)
        self.assertEqual(job["request_id"], request_id)
        self.assertEqual(job["source_type"], "demo")
        self.assertEqual(job["progress"], 100)
        self.assertGreaterEqual(len(job["logs"]), 5)
        self.assertEqual(len(job["stage_history"]), 5)
        self.assertTrue(job["illustration_meta"]["regenerated"])
        self.assertEqual(job["illustration_meta"]["regenerate_action"], "local_refined_prompt")
        self.assertEqual(job["intent"]["chart_type"], payload["pipeline"]["intent"]["chart_type"])

    def test_path_to_asset_url_maps_output_files(self):
        self.assertEqual(path_to_asset_url("outputs/demo.png"), "/assets/outputs/demo.png")

    def test_normalizers_accept_known_values(self):
        self.assertEqual(normalize_chart_theme("business"), "business")
        self.assertEqual(normalize_chart_theme("unknown"), "tech")
        self.assertEqual(normalize_chart_type_override("line"), "line")
        self.assertEqual(normalize_chart_type_override("auto"), "")
        self.assertEqual(normalize_illustration_style("tech"), "tech")
        self.assertEqual(normalize_image_model("wanx"), "wanx")

    def test_local_chart_recommendation_covers_milestone_two_intents(self):
        cases = [
            (
                "不同品牌手机销量存在明显差异。",
                ["brand", "sales"],
                [["A", 120], ["B", 90]],
                "comparison",
                "bar",
            ),
            (
                "2020年到2023年销售额持续增长。",
                ["year", "sales"],
                [[2020, 100], [2021, 150], [2022, 210]],
                "trend",
                "line",
            ),
            (
                "收入构成中软件业务占45%，硬件业务占35%。",
                ["category", "share"],
                [["software", 45], ["hardware", 35]],
                "composition",
                "pie",
            ),
            (
                "学生成绩主要集中在80到90分之间。",
                ["bucket", "count"],
                [["60-70", 5], ["70-80", 18], ["80-90", 32]],
                "distribution",
                "bar",
            ),
            (
                "广告投入越高，销售额通常越高。",
                ["ad_spend", "sales"],
                [[10, 120], [20, 180], [30, 260]],
                "correlation",
                "scatter",
            ),
        ]

        for text, columns, rows, expected_intent, expected_chart in cases:
            with self.subTest(expected_intent=expected_intent):
                recommendation = _recommend_chart_intent(text, columns, rows)
                self.assertEqual(recommendation["intent_category"], expected_intent)
                self.assertEqual(recommendation["chart_type"], expected_chart)
                self.assertGreaterEqual(recommendation["confidence"], 0.6)
                self.assertTrue(recommendation["signals"])
                self.assertIn("推荐", recommendation["reason"])

    def test_image_clients_require_api_keys_for_remote_models(self):
        settings = SimpleNamespace(
            wanx_api_key="",
            wanx_base_url="https://example.test/wanx",
            wanx_model="wanx-test",
            flux_api_key="",
            flux_base_url="https://example.test/flux",
            flux_model_endpoint="endpoint",
            image_generation_timeout_seconds=1,
            image_poll_interval_seconds=0.01,
        )
        with patch("backend.image_clients.get_settings", return_value=settings):
            with self.assertRaisesRegex(RuntimeError, "WANX_API_KEY"):
                generate_wanx_image("prompt", Path(tempfile.gettempdir()) / "wanx-no-key.png", api_key="")
            with self.assertRaisesRegex(RuntimeError, "FLUX_API_KEY"):
                generate_flux_image("prompt", Path(tempfile.gettempdir()) / "flux-no-key.png", api_key="")

    def test_wanx_client_downloads_returned_image(self):
        output_path = Path(tempfile.gettempdir()) / "mock-wanx-output.png"
        post_response = Mock(ok=True)
        post_response.json.return_value = {
            "output": {
                "choices": [
                    {
                        "message": {
                            "content": [{"image": "https://example.test/wanx.png"}],
                        }
                    }
                ]
            }
        }
        get_response = Mock(content=b"png-bytes")
        get_response.raise_for_status.return_value = None
        try:
            with patch("backend.image_clients.requests.post", return_value=post_response) as mocked_post:
                with patch("backend.image_clients.requests.get", return_value=get_response) as mocked_get:
                    result = generate_wanx_image("business prompt", output_path, api_key="test-key")
            self.assertEqual(result, output_path.with_suffix(".png"))
            self.assertEqual(result.read_bytes(), b"png-bytes")
            self.assertEqual(mocked_post.call_args.kwargs["headers"]["Authorization"], "Bearer test-key")
            self.assertEqual(mocked_get.call_args.args[0], "https://example.test/wanx.png")
        finally:
            output_path.with_suffix(".png").unlink(missing_ok=True)

    def test_flux_client_supports_direct_and_polled_results(self):
        output_direct = Path(tempfile.gettempdir()) / "mock-flux-direct.png"
        direct_response = Mock()
        direct_response.raise_for_status.return_value = None
        direct_response.json.return_value = {"result": {"sample": "https://example.test/flux-direct.png"}}
        direct_get_response = Mock(content=b"direct-bytes")
        direct_get_response.raise_for_status.return_value = None

        output_poll = Path(tempfile.gettempdir()) / "mock-flux-poll.png"
        submit_response = Mock()
        submit_response.raise_for_status.return_value = None
        submit_response.json.return_value = {"polling_url": "https://example.test/poll"}
        poll_response = Mock()
        poll_response.raise_for_status.return_value = None
        poll_response.json.return_value = {"status": "Ready", "result": "https://example.test/flux-ready.png"}
        poll_get_response = Mock(content=b"poll-bytes")
        poll_get_response.raise_for_status.return_value = None

        try:
            with patch("backend.image_clients.requests.post", return_value=direct_response):
                with patch("backend.image_clients.requests.get", return_value=direct_get_response):
                    direct_result = generate_flux_image("direct prompt", output_direct, api_key="flux-key")
            self.assertEqual(direct_result.read_bytes(), b"direct-bytes")

            with patch("backend.image_clients.requests.post", return_value=submit_response):
                with patch("backend.image_clients.requests.get", side_effect=[poll_response, poll_get_response]) as mocked_get:
                    poll_result = generate_flux_image("poll prompt", output_poll, api_key="flux-key")
            self.assertEqual(poll_result.read_bytes(), b"poll-bytes")
            self.assertEqual(mocked_get.call_args_list[0].args[0], "https://example.test/poll")
        finally:
            output_direct.with_suffix(".png").unlink(missing_ok=True)
            output_poll.with_suffix(".png").unlink(missing_ok=True)

    def test_flux_result_url_parser_accepts_supported_shapes(self):
        self.assertEqual(_resolve_flux_result_url({"sample": "https://example.test/a.png"}), "https://example.test/a.png")
        self.assertEqual(_resolve_flux_result_url({"result": {"image_url": "https://example.test/b.png"}}), "https://example.test/b.png")
        with self.assertRaisesRegex(RuntimeError, "downloadable image URL"):
            _resolve_flux_result_url({"result": {"status": "ready"}})


@unittest.skipUnless(Presentation is not None, "python-pptx is not installed")
class PptModuleTests(unittest.TestCase):
    def _build_sample_ppt(self) -> Path:
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(4.5), Inches(0.6)).text_frame.text = "Quarterly Revenue"
        slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(5.0), Inches(0.8)).text_frame.text = "Summary slide for chart insertion testing."
        slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(6.2), Inches(0.6), Inches(2.0), Inches(0.8)).text_frame.text = "Highlight"
        table = slide.shapes.add_table(4, 3, Inches(0.5), Inches(2.0), Inches(5.5), Inches(2.0)).table
        table.cell(0, 0).text = "month"
        table.cell(0, 1).text = "sales"
        table.cell(0, 2).text = "profit"
        table.cell(1, 0).text = "Jan"
        table.cell(1, 1).text = "120"
        table.cell(1, 2).text = "20"
        table.cell(2, 0).text = "Feb"
        table.cell(2, 1).text = "150"
        table.cell(2, 2).text = "35"
        table.cell(3, 0).text = "Mar"
        table.cell(3, 1).text = "180"
        table.cell(3, 2).text = "40"
        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
            presentation.save(tmp.name)
            return Path(tmp.name)

    def _build_merged_table_ppt(self) -> Path:
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        table = slide.shapes.add_table(4, 3, Inches(0.5), Inches(1.0), Inches(5.5), Inches(2.0)).table
        table.cell(0, 0).text = "Region"
        table.cell(0, 1).text = "Revenue"
        table.cell(0, 2).text = "Profit"
        table.cell(1, 0).text = "North"
        table.cell(2, 0).text = "South"
        table.cell(1, 1).text = "120"
        table.cell(1, 2).text = "20"
        table.cell(2, 1).text = "150"
        table.cell(2, 2).text = "35"
        table.cell(3, 0).text = "Total"
        table.cell(3, 1).text = "270"
        table.cell(3, 2).text = "55"
        table.cell(1, 0).merge(table.cell(2, 0))
        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
            presentation.save(tmp.name)
            return Path(tmp.name)

    def test_extract_slide_content_marks_empty_slide(self):
        presentation = Presentation()
        presentation.slides.add_slide(presentation.slide_layouts[6])
        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
            presentation.save(tmp.name)
            ppt_path = Path(tmp.name)
        try:
            parsed = extract_slide_content(ppt_path, 1)
            self.assertEqual(parsed.text_content, "")
            self.assertEqual(parsed.tables, [])
            self.assertTrue(parsed.diagnostics["is_empty"])
            self.assertEqual(parsed.diagnostics["picture_count"], 0)
            self.assertEqual(parsed.diagnostics["non_empty_text_shape_count"], 0)
        finally:
            ppt_path.unlink(missing_ok=True)

    def test_extract_slide_content_counts_picture_shapes(self):
        if Image is None:
            self.skipTest("Pillow is not installed")
        image_path = Path(tempfile.gettempdir()) / "codex-test-parser-picture.png"
        Image.new("RGB", (240, 160), "#d8ebff").save(image_path)
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        slide.shapes.add_picture(str(image_path), Inches(1), Inches(1), Inches(3), Inches(2))
        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
            presentation.save(tmp.name)
            ppt_path = Path(tmp.name)
        try:
            parsed = extract_slide_content(ppt_path, 1)
            self.assertEqual(parsed.diagnostics["picture_count"], 1)
            self.assertFalse(parsed.diagnostics["is_empty"])
            self.assertTrue(parsed.shapes[0]["has_picture"])
        finally:
            ppt_path.unlink(missing_ok=True)
            image_path.unlink(missing_ok=True)

    def test_extract_slide_content_uses_spatial_text_reading_order(self):
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        slide.shapes.add_textbox(Inches(5.0), Inches(0.8), Inches(2.0), Inches(0.5)).text_frame.text = "Top right"
        slide.shapes.add_textbox(Inches(0.8), Inches(0.8), Inches(2.0), Inches(0.5)).text_frame.text = "Top left"
        slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(2.0), Inches(0.5)).text_frame.text = "Lower left"
        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
            presentation.save(tmp.name)
            ppt_path = Path(tmp.name)
        try:
            parsed = extract_slide_content(ppt_path, 1)
            self.assertEqual(parsed.text_content.splitlines(), ["Top left", "Top right", "Lower left"])
            self.assertEqual(parsed.diagnostics["text_order"], [2, 1, 3])
        finally:
            ppt_path.unlink(missing_ok=True)

    def test_table_to_dataframe_resolves_merged_cells(self):
        ppt_path = self._build_merged_table_ppt()
        try:
            parsed = extract_slide_content(ppt_path, 1)
            dataframe = table_to_dataframe(Presentation(str(ppt_path)).slides[0].shapes[0].table)
            self.assertEqual(dataframe.columns.tolist(), ["Region", "Revenue", "Profit"])
            self.assertEqual(dataframe.iloc[0].to_dict(), {"Region": "North", "Revenue": "120", "Profit": "20"})
            self.assertEqual(dataframe.iloc[1].to_dict(), {"Region": "North", "Revenue": "150", "Profit": "35"})
            self.assertTrue(parsed.tables[0]["cell_matrix"][1][0]["is_merge_origin"])
            self.assertTrue(parsed.tables[0]["cell_matrix"][2][0]["is_spanned"])
        finally:
            ppt_path.unlink(missing_ok=True)

    def test_insert_chart_to_pptx_replaces_table_region(self):
        ppt_path = self._build_sample_ppt()
        output_chart = Path(tempfile.gettempdir()) / "insert_test_chart.png"
        output_ppt = Path(tempfile.gettempdir()) / "insert_test_output.pptx"
        try:
            parsed = extract_slide_content(ppt_path, 1)
            table = parsed.tables[0]
            records = [dict(zip(table["columns"], row)) for row in table["rows"]]
            chart = generate_chart(records, "bar", output_path=output_chart, title="Revenue Overview")
            result = insert_chart_to_pptx(ppt_path=ppt_path, chart_image_path=chart.output_path, slide_number=1, chart_title=chart.title, chart_spec=chart.to_dict(), shapes=parsed.shapes, output_path=output_ppt)
            self.assertTrue(Path(result.output_path).exists())
            self.assertTrue(result.replaced_table)
            enhanced = Presentation(result.output_path)
            self.assertEqual(len(enhanced.slides), 1)
            self.assertFalse(any(getattr(shape, "has_table", False) for shape in enhanced.slides[0].shapes))
        finally:
            ppt_path.unlink(missing_ok=True)
            output_chart.unlink(missing_ok=True)
            output_ppt.unlink(missing_ok=True)

    def test_choose_asset_regions_avoids_dense_default_chart_area(self):
        slide_width = int(Inches(10))
        slide_height = int(Inches(7.5))
        dense_shapes = [
            {
                "index": 1,
                "shape_type": "TEXT_BOX",
                "has_text": True,
                "has_table": False,
                "left": int(slide_width * 0.30),
                "top": int(slide_height * 0.18),
                "width": int(slide_width * 0.26),
                "height": int(slide_height * 0.48),
            },
            {
                "index": 2,
                "shape_type": "TEXT_BOX",
                "has_text": True,
                "has_table": False,
                "left": int(slide_width * 0.12),
                "top": int(slide_height * 0.18),
                "width": int(slide_width * 0.16),
                "height": int(slide_height * 0.18),
            },
        ]

        chart_region, illustration_region = _choose_asset_regions(slide_width, slide_height, dense_shapes)
        default_chart_region = (int(slide_width * 0.08), int(slide_height * 0.20), int(slide_width * 0.50), int(slide_height * 0.58))

        def total_overlap(region):
            return sum(
                _overlap_area(
                    region,
                    (shape["left"], shape["top"], shape["width"], shape["height"]),
                )
                for shape in dense_shapes
            )

        self.assertLess(total_overlap(chart_region), total_overlap(default_chart_region))
        self.assertEqual(_overlap_area(chart_region, illustration_region), 0)

    def test_choose_asset_regions_preserves_table_anchor_for_chart(self):
        slide_width = int(Inches(10))
        slide_height = int(Inches(7.5))
        table_anchor = (int(slide_width * 0.10), int(slide_height * 0.45), int(slide_width * 0.42), int(slide_height * 0.25))
        shapes = [
            {
                "index": 1,
                "shape_type": "TABLE",
                "has_text": False,
                "has_table": True,
                "left": table_anchor[0],
                "top": table_anchor[1],
                "width": table_anchor[2],
                "height": table_anchor[3],
            }
        ]

        chart_region, illustration_region = _choose_asset_regions(slide_width, slide_height, shapes, table_anchor, ignored_shape_index=1)

        self.assertEqual(chart_region, table_anchor)
        self.assertEqual(_overlap_area(chart_region, illustration_region), 0)

    def test_insert_generated_assets_appends_result_slide_for_dense_layout(self):
        if Image is None:
            self.skipTest("Pillow is not installed")
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        for row in range(4):
            for col in range(3):
                slide.shapes.add_textbox(
                    Inches(0.4 + col * 3.0),
                    Inches(1.0 + row * 1.25),
                    Inches(2.5),
                    Inches(0.85),
                ).text_frame.text = f"Dense content {row}-{col}"
        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
            presentation.save(tmp.name)
            ppt_path = Path(tmp.name)
        chart_path = Path(tempfile.gettempdir()) / "dense_layout_chart.png"
        illustration_path = Path(tempfile.gettempdir()) / "dense_layout_illustration.png"
        output_ppt = Path(tempfile.gettempdir()) / "dense_layout_output.pptx"
        try:
            Image.new("RGB", (800, 500), "#ffffff").save(chart_path)
            Image.new("RGB", (500, 400), "#ddeeff").save(illustration_path)
            parsed = extract_slide_content(ppt_path, 1)
            intent = {"chart_type": "bar", "semantic_mode": "local"}
            insert_generated_assets(
                ppt_path=ppt_path,
                output_path=output_ppt,
                slide_number=1,
                chart_path=chart_path,
                illustration_path=illustration_path,
                title="Dense layout result",
                subtitle="Dense page should be preserved.",
                intent=intent,
                shapes=parsed.shapes,
            )
            enhanced = Presentation(str(output_ppt))
            self.assertEqual(len(enhanced.slides), 2)
            self.assertEqual(intent["layout"]["insertion_mode"], "appendix")
            self.assertTrue(intent["layout"]["original_slide_preserved"])
            self.assertEqual(intent["layout"]["result_slide_number"], 2)
        finally:
            ppt_path.unlink(missing_ok=True)
            chart_path.unlink(missing_ok=True)
            illustration_path.unlink(missing_ok=True)
            output_ppt.unlink(missing_ok=True)

    def test_insert_generated_assets_accepts_manual_layout_override(self):
        if Image is None:
            self.skipTest("Pillow is not installed")
        presentation = Presentation()
        presentation.slides.add_slide(presentation.slide_layouts[6])
        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
            presentation.save(tmp.name)
            ppt_path = Path(tmp.name)
        chart_path = Path(tempfile.gettempdir()) / "manual_layout_chart.png"
        illustration_path = Path(tempfile.gettempdir()) / "manual_layout_illustration.png"
        output_ppt = Path(tempfile.gettempdir()) / "manual_layout_output.pptx"
        try:
            Image.new("RGB", (800, 450), "#ffffff").save(chart_path)
            Image.new("RGB", (800, 450), "#ddeeff").save(illustration_path)
            intent = {"chart_type": "bar", "semantic_mode": "local"}
            insert_generated_assets(
                ppt_path=ppt_path,
                output_path=output_ppt,
                slide_number=1,
                chart_path=chart_path,
                illustration_path=illustration_path,
                intent=intent,
                layout_override={
                    "chartX": 12,
                    "chartY": 18,
                    "chartScale": 40,
                    "illustrationX": 56,
                    "illustrationY": 24,
                    "illustrationScale": 30,
                },
            )
            enhanced = Presentation(str(output_ppt))
            self.assertEqual(len(enhanced.slides), 1)
            self.assertEqual(intent["layout"]["insertion_mode"], "manual_override")
            self.assertTrue(intent["layout"]["manual_override"])
            self.assertEqual(intent["layout"]["chart_region"]["left"], int(enhanced.slide_width * 0.12))
            self.assertEqual(intent["layout"]["illustration_region"]["left"], int(enhanced.slide_width * 0.56))
        finally:
            ppt_path.unlink(missing_ok=True)
            chart_path.unlink(missing_ok=True)
            illustration_path.unlink(missing_ok=True)
            output_ppt.unlink(missing_ok=True)

    def test_apply_batch_layout_overrides_writes_manual_layout_pptx(self):
        if Image is None:
            self.skipTest("Pillow is not installed")
        presentation = Presentation()
        presentation.slides.add_slide(presentation.slide_layouts[6])
        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
            presentation.save(tmp.name)
            ppt_path = Path(tmp.name)
        chart_path = Path(tempfile.gettempdir()) / "batch_manual_chart.png"
        illustration_path = Path(tempfile.gettempdir()) / "batch_manual_illustration.png"
        try:
            Image.new("RGB", (800, 450), "#ffffff").save(chart_path)
            Image.new("RGB", (800, 450), "#ddeeff").save(illustration_path)
            payload = apply_batch_layout_overrides(
                ppt_path,
                [
                    {
                        "slide_number": 1,
                        "chart_path": str(chart_path),
                        "illustration_path": str(illustration_path),
                        "intent": {"chart_type": "line", "semantic_mode": "local"},
                        "layout_override": {
                            "chartX": 8,
                            "chartY": 10,
                            "chartScale": 36,
                            "illustrationX": 52,
                            "illustrationY": 20,
                            "illustrationScale": 34,
                        },
                    }
                ],
                batch_request_id="test-layout",
            )
            self.assertEqual(payload["applied_count"], 1)
            self.assertTrue(Path(payload["final_pptx_path"]).exists())
            self.assertTrue(payload["final_pptx_url"].startswith("/assets/outputs/"))
            self.assertEqual(payload["slides"][0]["layout"]["insertion_mode"], "manual_override")
            Path(payload["final_pptx_path"]).unlink(missing_ok=True)
        finally:
            ppt_path.unlink(missing_ok=True)
            chart_path.unlink(missing_ok=True)
            illustration_path.unlink(missing_ok=True)


    def test_insert_chart_to_pptx_reuses_adjacent_title_and_legend_regions(self):
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(5.6), Inches(0.7)).text_frame.text = "Old title placeholder"
        table = slide.shapes.add_table(4, 3, Inches(0.5), Inches(1.5), Inches(5.5), Inches(2.0)).table
        table.cell(0, 0).text = "month"
        table.cell(0, 1).text = "sales"
        table.cell(0, 2).text = "profit"
        table.cell(1, 0).text = "Jan"
        table.cell(1, 1).text = "120"
        table.cell(1, 2).text = "20"
        table.cell(2, 0).text = "Feb"
        table.cell(2, 1).text = "150"
        table.cell(2, 2).text = "35"
        table.cell(3, 0).text = "Mar"
        table.cell(3, 1).text = "180"
        table.cell(3, 2).text = "40"
        slide.shapes.add_textbox(Inches(0.5), Inches(3.7), Inches(5.5), Inches(0.55)).text_frame.text = "Old footer placeholder"
        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
            presentation.save(tmp.name)
            ppt_path = Path(tmp.name)
        output_chart = Path(tempfile.gettempdir()) / "insert_layout_chart.png"
        output_ppt = Path(tempfile.gettempdir()) / "insert_layout_output.pptx"
        try:
            parsed = extract_slide_content(ppt_path, 1)
            records = [dict(zip(parsed.tables[0]["columns"], row)) for row in parsed.tables[0]["rows"]]
            chart = generate_chart(records, "bar", output_path=output_chart, title="Updated Revenue Overview")
            result = insert_chart_to_pptx(
                ppt_path=ppt_path,
                chart_image_path=chart.output_path,
                slide_number=1,
                chart_title=chart.title,
                chart_spec={"y_columns": ["sales", "profit"]},
                shapes=parsed.shapes,
                output_path=output_ppt,
            )
            enhanced = Presentation(result.output_path)
            texts = [shape.text.strip() for shape in enhanced.slides[0].shapes if getattr(shape, "has_text_frame", False)]
            self.assertIn("Updated Revenue Overview", texts)
            self.assertTrue(any("Legend: sales, profit" in text for text in texts))
            self.assertFalse(any("Old title placeholder" in text for text in texts))
            self.assertFalse(any("Old footer placeholder" in text for text in texts))
        finally:
            ppt_path.unlink(missing_ok=True)
            output_chart.unlink(missing_ok=True)
            output_ppt.unlink(missing_ok=True)

class ChartThemeTests(unittest.TestCase):
    def test_generate_chart_supports_all_milestone_two_chart_types(self):
        if Image is None:
            self.skipTest("Pillow is not installed")
        records = [
            {"label": "A", "value": 10, "other": 8, "third": 3},
            {"label": "B", "value": 18, "other": 12, "third": 5},
            {"label": "C", "value": 12, "other": 20, "third": 9},
            {"label": "D", "value": 24, "other": 16, "third": 7},
        ]
        chart_types = ["bar", "line", "pie", "scatter", "area", "histogram", "box", "heatmap"]
        output_paths = [Path(tempfile.gettempdir()) / f"m2_chart_type_{chart_type}.png" for chart_type in chart_types]
        try:
            for chart_type, output_path in zip(chart_types, output_paths):
                chart = generate_chart(records, chart_type, output_path=output_path, title=f"{chart_type} check")
                self.assertEqual(chart.chart_type, chart_type)
                self.assertFalse(chart.fallback)
                self.assertTrue(Path(chart.output_path).exists())
                self.assertGreater(Path(chart.output_path).stat().st_size, 0)
        finally:
            for output_path in output_paths:
                output_path.unlink(missing_ok=True)

    def test_generate_chart_renders_placeholder_for_empty_data(self):
        if Image is None:
            self.skipTest("Pillow is not installed")
        output_chart = Path(tempfile.gettempdir()) / "empty_chart_placeholder.png"
        try:
            chart = generate_chart([], "bar", output_path=output_chart, title="Empty Data")
            self.assertTrue(chart.fallback)
            self.assertEqual(chart.quality_status, "fallback")
            self.assertTrue(chart.review_required)
            self.assertIn("empty", " ".join(chart.warnings or []).lower())
            self.assertTrue(output_chart.exists())
            self.assertGreater(output_chart.stat().st_size, 0)
        finally:
            output_chart.unlink(missing_ok=True)

    def test_generate_chart_handles_missing_and_invalid_values(self):
        if Image is None:
            self.skipTest("Pillow is not installed")
        output_chart = Path(tempfile.gettempdir()) / "missing_values_chart.png"
        try:
            chart = generate_chart(
                [
                    {"label": "A", "value": "10"},
                    {"label": "B", "value": None},
                    {"label": "C", "value": "not-a-number"},
                ],
                "line",
                output_path=output_chart,
                title="Missing Values",
            )
            self.assertFalse(chart.fallback)
            self.assertGreater(chart.quality_score, 0)
            self.assertLess(chart.quality_checks["numeric_coverage"], 1)
            self.assertEqual(chart.quality_status, "review")
            self.assertTrue(chart.review_required)
            self.assertIn("value_labels", chart.render_notes)
            self.assertTrue(output_chart.exists())
        finally:
            output_chart.unlink(missing_ok=True)

    def test_generate_chart_reports_quality_metadata_and_sampling(self):
        if Image is None:
            self.skipTest("Pillow is not installed")
        output_chart = Path(tempfile.gettempdir()) / "sampled_quality_chart.png"
        records = [{"label": f"Segment {index}", "value": index * 10, "other": index * 3} for index in range(1, 15)]
        try:
            chart = generate_chart(records, "bar", output_path=output_chart, title="Quality Metadata")
            spec = chart.to_dict()
            self.assertFalse(chart.fallback)
            self.assertEqual(spec["data_points"], 10)
            self.assertEqual(spec["series_count"], 2)
            self.assertGreaterEqual(spec["quality_score"], 7)
            self.assertEqual(spec["quality_status"], "attention")
            self.assertFalse(spec["review_required"])
            self.assertEqual(spec["quality_checks"]["readability"], "sampled")
            self.assertTrue(any("sampled 10 of 14" in warning for warning in spec["warnings"]))
            self.assertIn("axis_ticks", spec["render_notes"])
            self.assertTrue(output_chart.exists())
        finally:
            output_chart.unlink(missing_ok=True)

    def test_generate_chart_marks_zero_baseline_for_negative_bars(self):
        if Image is None:
            self.skipTest("Pillow is not installed")
        output_chart = Path(tempfile.gettempdir()) / "negative_bar_quality.png"
        try:
            chart = generate_chart(
                [
                    {"label": "North", "value": 42},
                    {"label": "South", "value": -18},
                    {"label": "West", "value": 24},
                ],
                "bar",
                output_path=output_chart,
                title="Positive and Negative Values",
            )
            self.assertFalse(chart.fallback)
            self.assertIn("zero_baseline", chart.render_notes)
            self.assertEqual(chart.quality_checks["value_range"], ["-18", "42"])
            self.assertTrue(output_chart.exists())
        finally:
            output_chart.unlink(missing_ok=True)

    def test_generate_chart_groups_extra_pie_slices_into_other(self):
        if Image is None:
            self.skipTest("Pillow is not installed")
        output_chart = Path(tempfile.gettempdir()) / "pie_other_grouped.png"
        records = [{"label": f"Category {index}", "value": index} for index in range(1, 9)]
        try:
            chart = generate_chart(records, "pie", output_path=output_chart, title="Grouped Pie")
            self.assertFalse(chart.fallback)
            self.assertEqual(chart.data_points, 6)
            self.assertEqual(chart.quality_status, "attention")
            self.assertIn("pie_other_grouped", chart.render_notes)
            self.assertTrue(any("grouped 3 small slices into Other" in warning for warning in chart.warnings or []))
            self.assertTrue(output_chart.exists())
        finally:
            output_chart.unlink(missing_ok=True)

    def test_generate_chart_adds_synthetic_series_for_scatter_and_heatmap(self):
        if Image is None:
            self.skipTest("Pillow is not installed")
        scatter_path = Path(tempfile.gettempdir()) / "synthetic_scatter.png"
        heatmap_path = Path(tempfile.gettempdir()) / "synthetic_heatmap.png"
        records = [{"label": "A", "value": 3}, {"label": "B", "value": 5}]
        try:
            scatter = generate_chart(records, "scatter", output_path=scatter_path)
            heatmap = generate_chart(records, "heatmap", output_path=heatmap_path)
            self.assertFalse(scatter.fallback)
            self.assertIn("_point_index", scatter.y_columns)
            self.assertIn("scatter_synthetic_index", scatter.render_notes)
            self.assertTrue(any("synthetic" in warning.lower() for warning in scatter.warnings or []))
            self.assertFalse(heatmap.fallback)
            self.assertEqual(len(heatmap.y_columns), 2)
            self.assertTrue(any("baseline" in warning.lower() for warning in heatmap.warnings or []))
        finally:
            scatter_path.unlink(missing_ok=True)
            heatmap_path.unlink(missing_ok=True)

    def test_generate_chart_sanitizes_non_positive_pie_values(self):
        if Image is None:
            self.skipTest("Pillow is not installed")
        output_chart = Path(tempfile.gettempdir()) / "pie_sanitized.png"
        try:
            chart = generate_chart(
                [{"label": "A", "value": -5}, {"label": "B", "value": 0}],
                "pie",
                output_path=output_chart,
            )
            self.assertFalse(chart.fallback)
            self.assertTrue(any("placeholder slices" in warning for warning in chart.warnings or []))
            self.assertTrue(output_chart.exists())
        finally:
            output_chart.unlink(missing_ok=True)

    def test_generate_chart_uses_tech_theme_background(self):
        if Image is None:
            self.skipTest("Pillow is not installed")
        output_chart = Path(tempfile.gettempdir()) / "tech_theme_chart.png"
        try:
            chart = generate_chart(
                [
                    {"month": "Jan", "sales": 120},
                    {"month": "Feb", "sales": 150},
                    {"month": "Mar", "sales": 180},
                ],
                "bar",
                output_path=output_chart,
                title="Theme Check",
            )
            image = Image.open(chart.output_path)
            self.assertEqual(image.getpixel((10, 10)), (7, 17, 31))
        finally:
            output_chart.unlink(missing_ok=True)


    def test_generate_chart_supports_business_theme(self):
        if Image is None:
            self.skipTest("Pillow is not installed")
        output_chart = Path(tempfile.gettempdir()) / "business_theme_chart.png"
        try:
            chart = generate_chart(
                [
                    {"month": "Jan", "sales": 120},
                    {"month": "Feb", "sales": 150},
                    {"month": "Mar", "sales": 180},
                ],
                "bar",
                output_path=output_chart,
                title="Business Theme Check",
                theme="business",
            )
            image = Image.open(chart.output_path)
            self.assertEqual(chart.theme, "business")
            self.assertEqual(image.getpixel((10, 10)), (244, 247, 251))
        finally:
            output_chart.unlink(missing_ok=True)


class AppTests(unittest.TestCase):
    def test_create_app_registers_expected_routes(self):
        app = create_app()
        route_paths = {route.path for route in app.routes}
        self.assertIn("/api/health", route_paths)
        self.assertIn("/api/pipeline", route_paths)
        self.assertIn("/api/process", route_paths)
        self.assertIn("/api/process-batch", route_paths)
        self.assertIn("/api/batch-layout", route_paths)
        self.assertIn("/api/demo-chart", route_paths)
        self.assertIn("/api/slide-preview", route_paths)
        self.assertIn("/api/parse-slides", route_paths)
        self.assertIn("/api/auth/login", route_paths)
        self.assertIn("/api/jobs", route_paths)
        self.assertIn("/api/jobs/{request_id}", route_paths)

    def test_health_payload_exposes_image_provider_flags(self):
        payload = build_health_payload()
        self.assertIn("wanx_enabled", payload)
        self.assertIn("flux_enabled", payload)
        self.assertIn("batch-processing", payload["features"])
        self.assertIn("manual-layout-writeback", payload["features"])
        self.assertEqual(payload["database_engine"], "sqlite")

    def test_health_endpoint_returns_structured_payload(self):
        client = TestClient(create_app())
        response = client.get("/api/health")
        self.assertEqual(response.status_code, 200)
        payload = response.json()
        self.assertIn("batch-processing", payload["features"])
        self.assertIn("manual-layout-writeback", payload["features"])
        self.assertIsInstance(payload["database_enabled"], bool)

    def test_authenticate_or_create_user_persists_user(self):
        user = authenticate_or_create_user("codex-demo", "demo123")
        self.assertEqual(user["username"], "codex-demo")
        same_user = authenticate_or_create_user("codex-demo", "demo123")
        self.assertEqual(user["id"], same_user["id"])

    def test_job_detail_endpoint_returns_pipeline_metadata(self):
        payload = process_demo_text("Alpha\nBeta", illustration_style="auto", image_model="local")
        request_id = payload["pipeline"]["request_id"]
        client = TestClient(create_app())
        response = client.get(f"/api/jobs/{request_id}")
        self.assertEqual(response.status_code, 200)
        job = response.json()["job"]
        self.assertEqual(job["request_id"], request_id)
        self.assertIn("logs", job)
        self.assertIn("stage_history", job)
        self.assertTrue(job["illustration_meta"]["regenerated"])
        self.assertTrue(job["chart_image_url"].startswith("/assets/outputs/"))


if __name__ == "__main__":
    unittest.main()
