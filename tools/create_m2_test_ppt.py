from __future__ import annotations

from pathlib import Path
import tempfile

from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUTPUT = ROOT / "outputs" / "milestone2_manual_test_input.pptx"

WIDE_W = Inches(13.333)
WIDE_H = Inches(7.5)

COLORS = {
    "ink": RGBColor(24, 31, 42),
    "muted": RGBColor(92, 103, 120),
    "line": RGBColor(217, 224, 235),
    "paper": RGBColor(247, 249, 252),
    "white": RGBColor(255, 255, 255),
    "teal": RGBColor(0, 128, 128),
    "green": RGBColor(47, 125, 82),
    "coral": RGBColor(231, 84, 75),
    "gold": RGBColor(224, 157, 51),
    "violet": RGBColor(88, 76, 150),
    "navy": RGBColor(28, 44, 79),
}


def set_fill(shape, color: RGBColor) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = color


def set_text(
    shape,
    text: str,
    *,
    size: int = 16,
    color: RGBColor | None = None,
    bold: bool = False,
    align: PP_ALIGN | None = None,
) -> None:
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    if align is not None:
        paragraph.alignment = align
    for run in paragraph.runs:
        run.font.name = "Microsoft YaHei"
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color or COLORS["ink"]


def add_bg(slide, color: RGBColor = COLORS["paper"]) -> None:
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, WIDE_W, WIDE_H)
    set_fill(rect, color)
    rect.line.fill.background()


def add_title(slide, title: str, subtitle: str | None = None) -> None:
    box = slide.shapes.add_textbox(Inches(0.55), Inches(0.36), Inches(8.6), Inches(0.62))
    set_text(box, title, size=28, bold=True, color=COLORS["ink"])
    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.58), Inches(0.98), Inches(9.2), Inches(0.36))
        set_text(sub, subtitle, size=12, color=COLORS["muted"])


def add_footer(slide, label: str) -> None:
    box = slide.shapes.add_textbox(Inches(0.55), Inches(7.03), Inches(6.0), Inches(0.25))
    set_text(box, label, size=9, color=COLORS["muted"])


def add_card(slide, left, top, width, height, title: str, body: str, color: RGBColor) -> None:
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    set_fill(card, COLORS["white"])
    card.line.color.rgb = COLORS["line"]
    stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(0.09), height)
    set_fill(stripe, color)
    head = slide.shapes.add_textbox(left + Inches(0.24), top + Inches(0.18), width - Inches(0.4), Inches(0.28))
    set_text(head, title, size=14, bold=True, color=COLORS["ink"])
    txt = slide.shapes.add_textbox(left + Inches(0.24), top + Inches(0.58), width - Inches(0.4), height - Inches(0.78))
    set_text(txt, body, size=11, color=COLORS["muted"])


def add_table(slide, left, top, width, height, rows: list[list[str]], accent: RGBColor) -> None:
    table_shape = slide.shapes.add_table(len(rows), len(rows[0]), left, top, width, height)
    table = table_shape.table
    for row_index, row in enumerate(rows):
        for col_index, value in enumerate(row):
            cell = table.cell(row_index, col_index)
            cell.text = str(value)
            cell.margin_left = Inches(0.04)
            cell.margin_right = Inches(0.04)
            cell.margin_top = Inches(0.03)
            cell.margin_bottom = Inches(0.03)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER if row_index == 0 or col_index > 0 else PP_ALIGN.LEFT
                for run in paragraph.runs:
                    run.font.name = "Microsoft YaHei"
                    run.font.size = Pt(10 if len(rows) > 7 else 11)
                    run.font.bold = row_index == 0
                    run.font.color.rgb = COLORS["white"] if row_index == 0 else COLORS["ink"]
            cell.fill.solid()
            cell.fill.fore_color.rgb = accent if row_index == 0 else RGBColor(255, 255, 255)
    return table_shape


def add_metric(slide, left, top, value: str, label: str, color: RGBColor) -> None:
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, Inches(1.25), Inches(1.25))
    set_fill(circle, color)
    set_text(circle, value, size=20, bold=True, color=COLORS["white"], align=PP_ALIGN.CENTER)
    txt = slide.shapes.add_textbox(left - Inches(0.2), top + Inches(1.34), Inches(1.65), Inches(0.32))
    set_text(txt, label, size=10, color=COLORS["muted"], align=PP_ALIGN.CENTER)


def create_placeholder_image() -> Path:
    path = Path(tempfile.gettempdir()) / "m2_test_embedded_image.png"
    image = Image.new("RGB", (1200, 720), (238, 243, 248))
    draw = ImageDraw.Draw(image)
    for i in range(0, 1200, 24):
        color = (220 - i // 18, 236 - i // 28, 238)
        draw.line([(i, 0), (1200, min(720, i))], fill=color, width=10)
    draw.rounded_rectangle([110, 100, 1090, 620], radius=36, fill=(255, 255, 255), outline=(0, 128, 128), width=8)
    draw.ellipse([190, 180, 420, 410], fill=(231, 84, 75))
    draw.rectangle([500, 190, 980, 250], fill=(28, 44, 79))
    draw.rectangle([500, 310, 900, 360], fill=(0, 128, 128))
    draw.rectangle([500, 430, 1020, 480], fill=(224, 157, 51))
    image.save(path)
    return path


def build_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = WIDE_W
    prs.slide_height = WIDE_H
    blank = prs.slide_layouts[6]

    # 1. Cover and overall metadata.
    slide = prs.slides.add_slide(blank)
    add_bg(slide, COLORS["navy"])
    cover = slide.shapes.add_textbox(Inches(0.7), Inches(0.72), Inches(8.4), Inches(1.25))
    set_text(cover, "Milestone 2 测试 PPT", size=38, bold=True, color=COLORS["white"])
    sub = slide.shapes.add_textbox(Inches(0.75), Inches(1.88), Inches(8.8), Inches(0.8))
    set_text(sub, "用于验证 PPT 解析、图表推荐、配图生成、批量翻页预览与手动位置微调", size=17, color=RGBColor(219, 232, 244))
    add_metric(slide, Inches(8.7), Inches(0.85), "12", "测试页", COLORS["teal"])
    add_metric(slide, Inches(10.2), Inches(0.85), "5", "图表意图", COLORS["coral"])
    add_metric(slide, Inches(11.7), Inches(0.85), "3", "边界页", COLORS["gold"])
    add_card(slide, Inches(0.82), Inches(3.15), Inches(3.5), Inches(1.9), "验收目标", "生成增强版 PPT 后检查：图表是否贴合数据、配图是否去同质化、写回是否避开核心内容。", COLORS["teal"])
    add_card(slide, Inches(4.72), Inches(3.15), Inches(3.5), Inches(1.9), "推荐流程", "上传本文件，选择批量生成，随后逐页翻看结果并尝试调整图表和配图位置。", COLORS["coral"])
    add_card(slide, Inches(8.62), Inches(3.15), Inches(3.5), Inches(1.9), "判定标准", "不存在流程失败、PPT 无法打开、核心标题正文遮挡或明显主题错配。", COLORS["gold"])
    add_footer(slide, "M2 manual test input | cover")

    # 2. Trend.
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_title(slide, "趋势场景：季度活跃用户与收入增长", "期望推荐折线图或面积图，配图应表现增长、产品运营或用户留存。")
    add_table(slide, Inches(0.68), Inches(1.55), Inches(7.3), Inches(3.2), [
        ["季度", "活跃用户(万)", "收入(万元)", "留存率(%)"],
        ["2024 Q1", "28", "320", "61"],
        ["2024 Q2", "34", "410", "64"],
        ["2024 Q3", "46", "560", "68"],
        ["2024 Q4", "58", "760", "72"],
        ["2025 Q1", "71", "980", "75"],
        ["2025 Q2", "83", "1210", "78"],
    ], COLORS["teal"])
    add_card(slide, Inches(8.35), Inches(1.55), Inches(4.25), Inches(1.4), "解析检查", "表头应识别为季度、活跃用户、收入、留存率。", COLORS["teal"])
    add_card(slide, Inches(8.35), Inches(3.25), Inches(4.25), Inches(1.4), "人工验收重点", "图表需要体现随时间上升趋势，不能错误生成占比类图。", COLORS["green"])
    add_footer(slide, "M2.4 / M2.5 trend data")

    # 3. Composition.
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_title(slide, "构成场景：业务收入来源占比", "期望推荐饼图或环形图，低占比项目可合并为 Other。")
    add_table(slide, Inches(0.68), Inches(1.5), Inches(5.9), Inches(3.1), [
        ["收入来源", "金额(万元)", "占比(%)"],
        ["企业订阅", "860", "43"],
        ["行业解决方案", "520", "26"],
        ["培训与咨询", "270", "13.5"],
        ["插件市场", "190", "9.5"],
        ["其他", "160", "8"],
    ], COLORS["violet"])
    add_card(slide, Inches(7.0), Inches(1.5), Inches(2.55), Inches(1.45), "不要生成", "坐标轴密集的柱状图、无意义散点图。", COLORS["coral"])
    add_card(slide, Inches(9.85), Inches(1.5), Inches(2.55), Inches(1.45), "配图方向", "收入结构、业务组合、服务矩阵。", COLORS["violet"])
    add_card(slide, Inches(7.0), Inches(3.25), Inches(5.4), Inches(1.25), "人工验收重点", "占比合计应接近 100%，标签不能互相遮挡。", COLORS["gold"])
    add_footer(slide, "M2.4 / M2.5 composition data")

    # 4. Comparison.
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_title(slide, "对比场景：区域销售与满意度排名", "期望推荐柱状图，配图应体现区域网络或市场扩张。")
    add_table(slide, Inches(0.62), Inches(1.42), Inches(8.1), Inches(3.75), [
        ["区域", "销售额(万元)", "客户数", "满意度"],
        ["华东", "1280", "420", "4.7"],
        ["华南", "1060", "360", "4.5"],
        ["华北", "890", "310", "4.3"],
        ["西南", "730", "230", "4.1"],
        ["西北", "420", "140", "3.9"],
        ["东北", "390", "120", "3.8"],
    ], COLORS["green"])
    add_metric(slide, Inches(9.2), Inches(1.65), "1280", "最高销售额", COLORS["green"])
    add_metric(slide, Inches(10.95), Inches(1.65), "4.7", "最高满意度", COLORS["teal"])
    add_card(slide, Inches(9.05), Inches(3.65), Inches(3.45), Inches(1.25), "人工验收重点", "排序、轴标签和数值标注应清楚，不要遮挡原始表格。", COLORS["green"])
    add_footer(slide, "M2.4 / M2.5 comparison data")

    # 5. Correlation.
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_title(slide, "相关性场景：广告投入与转化收入", "期望推荐散点图，并尽量给出趋势线或相关性说明。")
    add_table(slide, Inches(0.65), Inches(1.48), Inches(6.7), Inches(3.9), [
        ["活动", "广告投入(万元)", "转化率(%)", "收入(万元)"],
        ["A", "12", "2.1", "78"],
        ["B", "18", "2.8", "130"],
        ["C", "24", "3.2", "176"],
        ["D", "31", "3.9", "236"],
        ["E", "38", "4.4", "302"],
        ["F", "45", "5.1", "389"],
    ], COLORS["coral"])
    add_card(slide, Inches(7.8), Inches(1.55), Inches(4.65), Inches(1.15), "解析检查", "广告投入和收入应作为两个连续数值轴。", COLORS["coral"])
    add_card(slide, Inches(7.8), Inches(2.95), Inches(4.65), Inches(1.15), "图表检查", "不要退化成只有活动名称的普通柱状图。", COLORS["violet"])
    add_card(slide, Inches(7.8), Inches(4.35), Inches(4.65), Inches(1.15), "配图检查", "适合营销增长、投放实验，不应出现错误的数据看板。", COLORS["gold"])
    add_footer(slide, "M2.4 / M2.5 correlation data")

    # 6. Distribution.
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_title(slide, "分布场景：任务处理耗时分布", "期望推荐直方图或箱线图，测试异常值和分布解释。")
    add_table(slide, Inches(0.6), Inches(1.45), Inches(7.0), Inches(4.35), [
        ["任务ID", "耗时(秒)", "文件大小(MB)", "页数"],
        ["T-01", "18", "3.2", "8"],
        ["T-02", "23", "4.1", "10"],
        ["T-03", "19", "2.8", "7"],
        ["T-04", "31", "5.4", "15"],
        ["T-05", "28", "4.8", "12"],
        ["T-06", "35", "6.1", "18"],
        ["T-07", "92", "11.5", "42"],
    ], COLORS["gold"])
    add_card(slide, Inches(8.05), Inches(1.62), Inches(4.35), Inches(1.55), "人工验收重点", "T-07 是明显异常值，图表需要保留它但不要让其他值完全不可读。", COLORS["gold"])
    add_card(slide, Inches(8.05), Inches(3.55), Inches(4.35), Inches(1.55), "配图方向", "系统监控、性能优化、任务队列。", COLORS["teal"])
    add_footer(slide, "M2.5 distribution and outlier data")

    # 7. Negative and missing values.
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_title(slide, "异常数据场景：利润有负值且存在缺失", "期望 fallback 稳定，不因 N/A、负数或零值中断流程。")
    add_table(slide, Inches(0.65), Inches(1.42), Inches(7.55), Inches(3.75), [
        ["月份", "收入(万元)", "成本(万元)", "利润(万元)"],
        ["1月", "210", "180", "30"],
        ["2月", "195", "205", "-10"],
        ["3月", "0", "65", "-65"],
        ["4月", "N/A", "140", "N/A"],
        ["5月", "260", "190", "70"],
        ["6月", "245", "260", "-15"],
    ], COLORS["coral"])
    add_card(slide, Inches(8.65), Inches(1.5), Inches(3.75), Inches(1.35), "图表检查", "负值应低于零基线，缺失值不能导致生成失败。", COLORS["coral"])
    add_card(slide, Inches(8.65), Inches(3.15), Inches(3.75), Inches(1.35), "版式检查", "新图表不要遮挡异常数据说明。", COLORS["violet"])
    add_footer(slide, "M2.5 fallback with negative and missing values")

    # 8. Dense text without table.
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_title(slide, "复杂文本场景：正文中隐藏指标", "期望解析出文本型数据，并给出合理图表或配图建议。")
    add_card(slide, Inches(0.65), Inches(1.5), Inches(5.85), Inches(4.4), "项目 A 复盘", "第一阶段完成率 62%，返工率 11%，平均响应时间 38 秒。第二阶段完成率提升到 78%，返工率下降到 7%，平均响应时间降至 24 秒。", COLORS["teal"])
    add_card(slide, Inches(6.85), Inches(1.5), Inches(5.85), Inches(4.4), "项目 B 复盘", "第一阶段满意度 4.1，第二阶段满意度 4.6。团队计划在下一轮把错误率控制在 5% 以下，并将处理耗时稳定在 20 秒以内。", COLORS["gold"])
    add_footer(slide, "M2.2 text-only numeric extraction")

    # 9. Real image slide.
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_title(slide, "图片页场景：含真实嵌入图片", "期望解析能识别图片元素，增强写回不要覆盖主图片。")
    image_path = create_placeholder_image()
    slide.shapes.add_picture(str(image_path), Inches(0.78), Inches(1.48), width=Inches(6.25), height=Inches(3.75))
    add_card(slide, Inches(7.42), Inches(1.5), Inches(4.65), Inches(1.45), "图片说明", "这是一张嵌入式图片，用于检查 picture_count 和写回避让。", COLORS["teal"])
    add_card(slide, Inches(7.42), Inches(3.25), Inches(4.65), Inches(1.45), "人工验收重点", "生成的图表或配图不能直接压在主图片上。", COLORS["coral"])
    add_footer(slide, "M2.2 picture detection / M2.3 layout avoidance")

    # 10. Intent conflict.
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_title(slide, "混合意图场景：趋势与构成同时出现", "期望系统优先选择主任务，人工检查推荐解释是否说得通。")
    add_table(slide, Inches(0.65), Inches(1.48), Inches(5.8), Inches(3.0), [
        ["季度", "总收入", "订阅占比", "服务占比"],
        ["Q1", "520", "52", "31"],
        ["Q2", "680", "55", "30"],
        ["Q3", "810", "57", "29"],
        ["Q4", "960", "60", "27"],
    ], COLORS["violet"])
    add_card(slide, Inches(6.95), Inches(1.48), Inches(5.35), Inches(1.25), "主问题", "如果用户关心收入增长，应优先趋势图。", COLORS["teal"])
    add_card(slide, Inches(6.95), Inches(3.05), Inches(5.35), Inches(1.25), "次问题", "如果用户关心收入结构，则构成图也合理。", COLORS["violet"])
    add_card(slide, Inches(6.95), Inches(4.62), Inches(5.35), Inches(1.1), "人工验收重点", "前端展示的推荐理由需要能解释选择。", COLORS["gold"])
    add_footer(slide, "M2.4 recommendation explanation conflict case")

    # 11. Dense layout.
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_title(slide, "版式压力场景：内容拥挤但核心区域不可遮挡", "期望写回算法选择安全区域，必要时用户可手动微调。")
    for i, (title, body, color) in enumerate([
        ("风险 1", "标题区和左侧正文是核心内容，不能被覆盖。", COLORS["coral"]),
        ("风险 2", "右下角预留较小区域，适合放置配图。", COLORS["gold"]),
        ("风险 3", "原始表格区域可作为图表替换候选。", COLORS["teal"]),
        ("风险 4", "导出微调版 PPT 后需要检查坐标是否生效。", COLORS["violet"]),
    ]):
        x = Inches(0.65 + (i % 2) * 3.2)
        y = Inches(1.52 + (i // 2) * 1.72)
        add_card(slide, x, y, Inches(2.85), Inches(1.35), title, body, color)
    add_table(slide, Inches(7.15), Inches(1.48), Inches(5.0), Inches(2.3), [
        ["区域", "指标", "状态"],
        ["左上", "标题", "不可遮挡"],
        ["左中", "正文卡片", "不可遮挡"],
        ["右上", "表格", "可替换"],
        ["右下", "空白", "可插入"],
    ], COLORS["green"])
    note = slide.shapes.add_textbox(Inches(7.15), Inches(4.18), Inches(5.0), Inches(0.95))
    set_text(note, "验收动作：批量生成后进入前端翻页，拖动或调整图表/配图位置，再导出微调版 PPT。", size=13, color=COLORS["ink"], bold=True)
    add_footer(slide, "M2.3 manual layout override stress case")

    # 12. Empty slide.
    slide = prs.slides.add_slide(blank)

    return prs


def main() -> None:
    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    presentation = build_presentation()
    presentation.save(OUTPUT)
    print(OUTPUT)


if __name__ == "__main__":
    main()
