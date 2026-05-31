from __future__ import annotations

import json
import sys
import time
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

sys.path.insert(0, str(Path(__file__).resolve().parents[1]))

from backend.services import process_local_ppt_batch


OUTPUT_DIR = Path("outputs")
SOURCE_PPT = OUTPUT_DIR / "week7_50_slide_smoke_source.pptx"
SUMMARY_JSON = OUTPUT_DIR / "week7_50_slide_smoke_summary.json"


def build_source_ppt(path: Path, slide_count: int = 50) -> Path:
    path.parent.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    for slide_number in range(1, slide_count + 1):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(8.5), Inches(0.6)).text_frame.text = (
            f"第 {slide_number} 页营收趋势测试"
        )
        table = slide.shapes.add_table(4, 2, Inches(0.7), Inches(1.35), Inches(5.8), Inches(2.2)).table
        table.cell(0, 0).text = "季度"
        table.cell(0, 1).text = "营收"
        table.cell(1, 0).text = "Q1"
        table.cell(1, 1).text = str(100 + slide_number)
        table.cell(2, 0).text = "Q2"
        table.cell(2, 1).text = str(130 + slide_number)
        table.cell(3, 0).text = "Q3"
        table.cell(3, 1).text = str(160 + slide_number)
    prs.save(path)
    return path


def main() -> None:
    start = time.perf_counter()
    source = build_source_ppt(SOURCE_PPT)
    result = process_local_ppt_batch(
        source,
        list(range(1, 51)),
        semantic_mode="local",
        chart_type_override="line",
        chart_theme="business",
        illustration_style="tech",
        image_model="local",
    )
    elapsed = time.perf_counter() - start
    final_pptx_path = Path(result["final_pptx_path"])
    completed = [
        slide
        for slide in result["slides"]
        if slide.get("status") == "completed" and Path(slide.get("chart_image", "")).exists()
    ]
    clip_scores = [
        float(slide.get("illustration_meta", {}).get("clip_score", 0))
        for slide in result["slides"]
    ]
    summary = {
        "source_ppt": str(source),
        "final_pptx_path": str(final_pptx_path),
        "requested_slides": 50,
        "processed_count": result["processed_count"],
        "completed_count": len(completed),
        "final_pptx_exists": final_pptx_path.exists(),
        "average_clip_score": round(sum(clip_scores) / len(clip_scores), 2) if clip_scores else 0,
        "elapsed_seconds": round(elapsed, 2),
        "pass": result["processed_count"] == 50 and len(completed) == 50 and final_pptx_path.exists(),
    }
    SUMMARY_JSON.write_text(json.dumps(summary, ensure_ascii=False, indent=2), encoding="utf-8")
    print(json.dumps(summary, ensure_ascii=False, indent=2))
    if not summary["pass"]:
        raise SystemExit(1)


if __name__ == "__main__":
    main()
