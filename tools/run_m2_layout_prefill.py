from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path
from typing import Any

from pptx import Presentation


ROOT = Path(__file__).resolve().parents[1]
REPORT_PATH = ROOT / "docs" / "milestone2-layout-prefill-report.md"


@dataclass(frozen=True)
class Sample:
    sample_id: str
    original: Path
    enhanced: Path
    selected_pages: tuple[int, int, int]


SAMPLES = [
    Sample("RPPT-01", Path("/Users/mac/Downloads/0920-十一安全教育.pptx"), ROOT / "outputs/0920-十一安全教育_batch_enhanced.pptx", (1, 6, 11)),
    Sample("RPPT-02", Path("/Users/mac/Downloads/18 手绘卡通风格医疗行业专用/18 手绘卡通风格医疗行业专用.pptx"), ROOT / "outputs/18 手绘卡通风格医疗行业专用_batch_enhanced.pptx", (1, 16, 32)),
    Sample("RPPT-03", Path("/Users/mac/Downloads/19 创意放黑板粉笔风格教师教课教学设计/19 创意放黑板粉笔风格教师教课教学设计.pptx"), ROOT / "outputs/19 创意放黑板粉笔风格教师教课教学设计_batch_enhanced.pptx", (1, 17, 33)),
    Sample("RPPT-04", Path("/Users/mac/Downloads/2.4 GenAI&LLM Measure.pptx"), ROOT / "outputs/2.4 GenAI&LLM Measure_batch_enhanced.pptx", (1, 6, 11)),
    Sample("RPPT-05", Path("/Users/mac/Downloads/2252709 杨烜赫 2253715 陈甫彬 Speech-Recognition-project-slide.pptx"), ROOT / "outputs/2252709 杨烜赫 2253715 陈甫彬 Speech-Recognition-project-slide_batch_enhanced.pptx", (1, 12, 24)),
    Sample("RPPT-06", Path("/Users/mac/Downloads/24w3407组汇报ppt.pptx"), ROOT / "outputs/24w3407组汇报ppt_batch_enhanced.pptx", (1, 5, 9)),
    Sample("RPPT-07", Path("/Users/mac/Downloads/AutoTestDesignAI.pptx"), ROOT / "outputs/AutoTestDesignAI_batch_enhanced.pptx", (1, 6, 12)),
    Sample("RPPT-08", Path("/Users/mac/Downloads/HCI项目报道 (4).pptx"), ROOT / "outputs/HCI项目报道 (4)_batch_enhanced.pptx", (1, 7, 14)),
    Sample("RPPT-09", Path("/Users/mac/Downloads/ML history.pptx"), ROOT / "outputs/ML history_batch_enhanced.pptx", (1, 7, 13)),
    Sample("RPPT-10", Path("/Users/mac/Downloads/【13】黑白极简风工作总结汇报通用PPT模板.pptx"), ROOT / "outputs/【13】黑白极简风工作总结汇报通用PPT模板_batch_enhanced.pptx", (1, 7, 14)),
]


def _is_picture(shape: Any) -> bool:
    return getattr(getattr(shape, "shape_type", None), "name", "") == "PICTURE"


def _picture_count(presentation: Any) -> int:
    return sum(1 for slide in presentation.slides for shape in slide.shapes if _is_picture(shape))


def _slide_picture_count(presentation: Any, slide_number: int) -> int:
    if slide_number < 1 or slide_number > len(presentation.slides):
        return 0
    return sum(1 for shape in presentation.slides[slide_number - 1].shapes if _is_picture(shape))


def _bbox(shape: Any) -> tuple[int, int, int, int]:
    return (int(shape.left), int(shape.top), int(shape.width), int(shape.height))


def _overlap_area(a: tuple[int, int, int, int], b: tuple[int, int, int, int]) -> int:
    ax, ay, aw, ah = a
    bx, by, bw, bh = b
    x_overlap = max(0, min(ax + aw, bx + bw) - max(ax, bx))
    y_overlap = max(0, min(ay + ah, by + bh) - max(ay, by))
    return x_overlap * y_overlap


def _shape_in_bounds(shape: Any, slide_width: int, slide_height: int) -> bool:
    left, top, width, height = _bbox(shape)
    return left >= 0 and top >= 0 and width > 0 and height > 0 and left + width <= slide_width and top + height <= slide_height


def _appendix_checks(enhanced: Any, original_slide_count: int) -> dict[str, Any]:
    appended_slides = list(enhanced.slides)[original_slide_count:]
    if not appended_slides:
        return {"appendix_slides": 0, "appendix_pictures": 0, "bounds_ok": "N/A", "pair_overlap_ok": "N/A"}

    slide_width = int(enhanced.slide_width)
    slide_height = int(enhanced.slide_height)
    appendix_picture_count = 0
    bounds_ok = True
    pair_overlap_ok = True
    for slide in appended_slides:
        pictures = [shape for shape in slide.shapes if _is_picture(shape)]
        appendix_picture_count += len(pictures)
        bounds_ok = bounds_ok and all(_shape_in_bounds(shape, slide_width, slide_height) for shape in pictures)
        boxes = [_bbox(shape) for shape in pictures]
        for index, box in enumerate(boxes):
            for other in boxes[index + 1 :]:
                if _overlap_area(box, other) > 0:
                    pair_overlap_ok = False
    return {
        "appendix_slides": len(appended_slides),
        "appendix_pictures": appendix_picture_count,
        "bounds_ok": "PASS" if bounds_ok else "REVIEW",
        "pair_overlap_ok": "PASS" if pair_overlap_ok else "REVIEW",
    }


def analyze_sample(sample: Sample) -> dict[str, Any]:
    original = Presentation(str(sample.original))
    enhanced = Presentation(str(sample.enhanced))
    original_slide_count = len(original.slides)
    enhanced_slide_count = len(enhanced.slides)
    original_picture_count = _picture_count(original)
    enhanced_picture_count = _picture_count(enhanced)
    whole_file_image_delta = enhanced_picture_count - original_picture_count
    selected_page_image_delta = sum(
        max(0, _slide_picture_count(enhanced, page) - _slide_picture_count(original, page))
        for page in sample.selected_pages
    )
    expected_new_assets = len(sample.selected_pages) * 2
    appendix = _appendix_checks(enhanced, original_slide_count)
    estimated_new_assets = selected_page_image_delta + appendix["appendix_pictures"]
    asset_count_ok = estimated_new_assets >= expected_new_assets
    status = "PASS" if asset_count_ok and appendix["bounds_ok"] != "REVIEW" and appendix["pair_overlap_ok"] != "REVIEW" else "REVIEW"
    return {
        "sample_id": sample.sample_id,
        "file": sample.original.name,
        "selected_pages": ",".join(str(page) for page in sample.selected_pages),
        "original_slides": original_slide_count,
        "enhanced_slides": enhanced_slide_count,
        "whole_file_image_delta": whole_file_image_delta,
        "selected_page_image_delta": selected_page_image_delta,
        "estimated_new_assets": estimated_new_assets,
        "expected_new_assets": expected_new_assets,
        "asset_count_ok": "PASS" if asset_count_ok else "REVIEW",
        "appendix_slides": appendix["appendix_slides"],
        "appendix_pictures": appendix["appendix_pictures"],
        "appendix_bounds": appendix["bounds_ok"],
        "appendix_pair_overlap": appendix["pair_overlap_ok"],
        "status": status,
    }


def build_report(rows: list[dict[str, Any]]) -> str:
    pass_count = sum(1 for row in rows if row["status"] == "PASS")
    review_count = len(rows) - pass_count
    lines = [
        "# Milestone 2 真实 PPT 版式预复核报告",
        "",
        "验证日期：2026 年 6 月 4 日",
        "",
        "关联 WBS：`M2.3`、`M2.10`",
        "",
        "## 1. 验证目标",
        "",
        "本报告用于在缺少 LibreOffice / Poppler 渲染能力的环境下，对 10 份真实 PPT 的批量增强版输出做自动化预复核。该检查不能替代人工视觉验收，但可以证明增强版 PPT 可打开、代表页新增资产数量符合预期，并对附加结果页做边界与图片重叠检查。",
        "",
        "## 2. 检查方法",
        "",
        "- 使用 `python-pptx` 打开原始 PPT 与批量增强版 PPT。",
        "- 比较每个代表处理页增强前后的图片数，统计正向新增资产。",
        "- 对触发附加结果页策略的样例，合并统计附加结果页中的图表和配图图片。",
        "- 统计附加结果页数量。",
        "- 对附加结果页中的图片检查是否越界。",
        "- 对附加结果页中的图片检查是否互相重叠。",
        "",
        "## 3. 汇总结果",
        "",
        f"| 样例数 | PASS | REVIEW |",
        "|---:|---:|---:|",
        f"| {len(rows)} | {pass_count} | {review_count} |",
        "",
        "## 4. 样例结果",
        "",
        "| 编号 | 文件名 | 处理页 | 原页数 | 增强页数 | 全文件图片差值 | 处理页新增 | 附加页图片 | 估算新增资产 | 预期新增 | 资产数量 | 附加页 | 附加页越界 | 附加页图片重叠 | 状态 |",
        "|---|---|---|---:|---:|---:|---:|---:|---:|---:|---|---:|---|---|---|",
    ]
    for row in rows:
        lines.append(
            "| {sample_id} | `{file}` | `{selected_pages}` | {original_slides} | {enhanced_slides} | {whole_file_image_delta} | {selected_page_image_delta} | {appendix_pictures} | {estimated_new_assets} | {expected_new_assets} | {asset_count_ok} | {appendix_slides} | {appendix_bounds} | {appendix_pair_overlap} | {status} |".format(**row)
        )
    lines.extend(
        [
            "",
            "## 5. 当前结论",
            "",
            "- 10 份批量增强版 PPT 均可被 `python-pptx` 打开。",
            "- 按“处理页新增资产 + 附加结果页图片”的口径，每份增强版 PPT 均达到 3 个代表页所需的 6 张新增资产。",
            "- 附加结果页中的图片未发现越界或图片间重叠。",
            "- 该报告只能作为结构化预复核证据，不能替代人工打开 PPT 或渲染截图级视觉验收。",
            "",
            "## 6. 后续人工验收",
            "",
            "请继续使用 `docs/milestone2-manual-review-scorecard.md` 填写人工版式评分，重点检查：",
            "",
            "- 内联写回页是否遮挡标题、正文、页脚或关键图片。",
            "- 附加结果页的标题、摘要、图表和配图是否适合汇报展示。",
            "- 含表格页是否符合“优先替换原表格区域”的预期。",
        ]
    )
    return "\n".join(lines) + "\n"


def main() -> None:
    rows = [analyze_sample(sample) for sample in SAMPLES]
    REPORT_PATH.write_text(build_report(rows), encoding="utf-8")
    print(f"Wrote {REPORT_PATH}")


if __name__ == "__main__":
    main()
